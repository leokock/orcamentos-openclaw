#!/usr/bin/env python3.11
# -*- coding: utf-8 -*-
"""
Gerador de Apresentação PowerPoint - Orçamento Paramétrico Edifício Rozzo
Cartesian Engenharia
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Cores Cartesian
NAVY = RGBColor(44, 62, 80)          # #2C3E50
BLUE = RGBColor(52, 152, 219)        # #3498DB
ORANGE = RGBColor(230, 126, 34)      # #E67E22
WHITE = RGBColor(255, 255, 255)      # #FFFFFF
LIGHT_GRAY = RGBColor(236, 240, 241) # #ECF0F1
GRAY = RGBColor(189, 195, 199)       # #BDC3C7

# Dados do projeto
PROJETO_NOME = "Edifício Rozzo — Vitório Demarche"
CLIENTE = "Rozzo Empreendimentos"
LOCAL = "Brusque/SC"
AC = 14854.30
UR = 115
UC = 3
PAVIMENTOS = 30
PAVIMENTOS_TIPO = 23
ELEVADORES = 3
VAGAS = 140
TERRENO = 1169.18
PRAZO = 36
CUB_ATUAL = 3050.00

# Custos
CUSTOS_RSM2 = [
    ("Supraestrutura", 722.66),
    ("Gerenciamento", 407.07),
    ("Instalações", 366.96),
    ("Esquadrias", 367.23),
    ("Infraestrutura", 230.40),
    ("Complementares", 214.32),
    ("Pisos", 193.69),
    ("Sist. Especiais", 193.39),
    ("Rev. Int. Parede", 171.85),
    ("Fachada", 170.21),
    ("Alvenaria", 148.70),
    ("Pintura", 133.12),
    ("Imprevistos", 97.79),
    ("Teto", 68.65),
    ("Impermeabilização", 56.43),
    ("Mov. Terra", 12.47),
    ("Climatização", 0.00),
    ("Louças e Metais", 0.00),
]

TOTAL_RSM2 = sum(c[1] for c in CUSTOS_RSM2)
TOTAL_VALOR = AC * TOTAL_RSM2
CUB_RATIO = TOTAL_RSM2 / CUB_ATUAL

# Benchmark
BENCHMARK = [
    ("Maison Beach", "Floripa", 12880, 25, 36, 3291, 1.20),
    ("Catena", "Floripa", 9242, 20, 30, 3376, 1.12),
    ("Connect", "Itajaí", 13144, 28, 36, 3150, 1.09),
    ("Dlohn", "Blumenau", 8500, 18, 24, 3080, 1.10),
    ("Eternity", "Itajaí", 11200, 22, 30, 3420, 1.20),
    ("Lorenzo", "Itajaí", 9800, 20, 28, 3190, 1.15),
]

def format_money(value):
    """Formata valores monetários em pt-BR"""
    return f"R$ {value:,.0f}".replace(",", ".")

def format_money_decimal(value):
    """Formata valores monetários com decimais em pt-BR"""
    return f"R$ {value:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

def format_percent(value):
    """Formata percentuais"""
    return f"{value:.1f}%"

def add_header_bar(slide, title_text):
    """Adiciona barra de título com accent laranja"""
    # Barra laranja
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.15), Inches(0.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ORANGE
    accent_bar.line.fill.background()
    
    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.25), Inches(0.05),
        Inches(9.5), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.text = title_text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = 'Calibri'

def add_kpi_card(slide, left, top, width, height, label, value, color=NAVY):
    """Adiciona card de KPI com borda colorida"""
    # Borda/fundo
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(3)
    
    # Label
    label_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.1),
        Inches(width - 0.2), Inches(0.3)
    )
    tf = label_box.text_frame
    tf.text = label
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Valor
    value_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.4),
        Inches(width - 0.2), Inches(height - 0.5)
    )
    tf = value_box.text_frame
    tf.text = value
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_presentation():
    """Cria apresentação completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 widescreen
    prs.slide_height = Inches(5.625)
    
    # Slide 1 - Capa
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background branco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # CARTESIAN ENGENHARIA (topo)
    title_top = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.4))
    tf = title_top.text_frame
    tf.text = "CARTESIAN ENGENHARIA"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Orçamento Paramétrico
    title_main = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.6))
    tf = title_main.text_frame
    tf.text = "Orçamento Paramétrico"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Nome do projeto
    project_name = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(0.5))
    tf = project_name.text_frame
    tf.text = PROJETO_NOME
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Cliente e local
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(8), Inches(0.4))
    tf = subtitle.text_frame
    tf.text = f"{CLIENTE} | {LOCAL}"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Data
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.4))
    tf = date_box.text_frame
    tf.text = "Março 2026"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2 - O que é o Orçamento Paramétrico
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "O que é o Orçamento Paramétrico")
    
    # Fórmula
    formula_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.6))
    tf = formula_box.text_frame
    tf.text = "Base Calibrada × CUB × Briefing = Custo Estimado"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Bullets
    bullets = [
        "✓ Rápido — Estimativa em minutos com base em 18 projetos reais",
        "✓ Ajustável — 25 parâmetros de briefing que recalculam automaticamente",
        "✓ Calibrado — Validado com orçamentos executivos da Cartesian"
    ]
    
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(2))
    tf = text_box.text_frame
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = NAVY
        p.font.name = 'Calibri'
        p.space_after = Pt(12)
    
    # Nota inferior
    note_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.5))
    tf = note_box.text_frame
    tf.text = "Não substitui o executivo — ANTECIPA decisões"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 3 - Dados do Empreendimento
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Dados do Empreendimento")
    
    # Grid de KPIs (3x2)
    kpis = [
        ("Área Construída", f"{AC:,.2f} m²".replace(",", "X").replace(".", ",").replace("X", ".")),
        ("Unidades", f"{UR} UR + {UC} UC"),
        ("Pavimentos", f"{PAVIMENTOS} total ({PAVIMENTOS_TIPO} tipo)"),
        ("Elevadores", str(ELEVADORES)),
        ("Vagas", f"~{VAGAS}"),
        ("Terreno", f"{TERRENO:,.2f} m²".replace(",", "X").replace(".", ",").replace("X", "."))
    ]
    
    positions = [
        (0.5, 1.2), (3.5, 1.2), (6.5, 1.2),
        (0.5, 3.2), (3.5, 3.2), (6.5, 3.2)
    ]
    
    for (label, value), (left, top) in zip(kpis, positions):
        add_kpi_card(slide, left, top, 2.5, 1.5, label, value, BLUE)
    
    # Slide 4 - Premissas (Briefing)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Premissas do Briefing")
    
    # Tabela 2 colunas
    premissas = [
        ("Tipo de Laje", "Cubetas"),
        ("Padrão", "Standard"),
        ("Fundação", "Hélice Contínua"),
        ("Prazo", "36 meses"),
        ("Contenção", "Sem"),
        ("Fachada", "Textura + pintura"),
        ("Gerador", "Com"),
        ("Pressurização", "Sem"),
    ]
    
    # Coluna 1
    y = 1.2
    for label, valor in premissas[:4]:
        # Label
        label_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(2), Inches(0.35))
        tf = label_box.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        # Valor
        value_box = slide.shapes.add_textbox(Inches(2.9), Inches(y), Inches(1.5), Inches(0.35))
        tf = value_box.text_frame
        tf.text = valor
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        y += 0.5
    
    # Coluna 2
    y = 1.2
    for label, valor in premissas[4:]:
        # Label
        label_box = slide.shapes.add_textbox(Inches(5.2), Inches(y), Inches(2), Inches(0.35))
        tf = label_box.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        # Valor
        value_box = slide.shapes.add_textbox(Inches(7.3), Inches(y), Inches(1.8), Inches(0.35))
        tf = value_box.text_frame
        tf.text = valor
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        y += 0.5
    
    # Nota
    note_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.5))
    tf = note_box.text_frame
    tf.text = "Premissas ajustáveis — mudanças recalculam automaticamente na planilha"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 5 - Resumo Executivo (KPIs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Resumo Executivo")
    
    # KPIs principais (2x2)
    kpis_exec = [
        ("Custo Total", format_money(TOTAL_VALOR), ORANGE),
        ("Custo por m²", format_money_decimal(TOTAL_RSM2), BLUE),
        ("CUB Ratio", f"{CUB_RATIO:.3f}", NAVY),
        ("Classificação", "Econômico-Médio", BLUE),
    ]
    
    positions = [
        (0.5, 1.2), (5.25, 1.2),
        (0.5, 3), (5.25, 3)
    ]
    
    for (label, value, color), (left, top) in zip(kpis_exec, positions):
        add_kpi_card(slide, left, top, 4.25, 1.5, label, value, color)
    
    # Base CUB
    note_box = slide.shapes.add_textbox(Inches(2), Inches(4.9), Inches(6), Inches(0.5))
    tf = note_box.text_frame
    tf.text = f"Base: CUB {format_money_decimal(CUB_ATUAL)} (mar/2026) • Valores sujeitos a ajuste conforme briefing"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 6 - Distribuição por Macrogrupo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Distribuição por Macrogrupo")
    
    # Ordenar por valor
    custos_sorted = sorted([(nome, valor) for nome, valor in CUSTOS_RSM2 if valor > 0], 
                          key=lambda x: x[1], reverse=True)
    
    y = 1
    max_width = 7
    for i, (nome, valor_rsm2) in enumerate(custos_sorted):
        perc = (valor_rsm2 / TOTAL_RSM2) * 100
        bar_width = (valor_rsm2 / custos_sorted[0][1]) * max_width
        
        # Nome
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2), Inches(0.18))
        tf = label_box.text_frame
        tf.text = nome
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        # Barra
        color = ORANGE if i < 5 else BLUE
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2.6), Inches(y),
            Inches(bar_width), Inches(0.18)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        # Valor e %
        value_box = slide.shapes.add_textbox(Inches(2.6 + bar_width + 0.1), Inches(y), Inches(1.5), Inches(0.18))
        tf = value_box.text_frame
        tf.text = f"{format_money_decimal(valor_rsm2)} ({perc:.1f}%)"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        y += 0.22
    
    # Slide 7 - Top 5 Custos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Top 5 Macrogrupos por Custo")
    
    # Header da tabela
    headers = ["Macrogrupo", "R$/m²", "% Total", "Valor Total"]
    x_positions = [0.5, 4, 6, 8]
    
    y = 1.2
    for header, x in zip(headers, x_positions):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.3))
        tf = box.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Background header
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(1.8), Inches(0.3)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        
    y = 1.6
    for i, (nome, valor_rsm2) in enumerate(custos_sorted[:5]):
        perc = (valor_rsm2 / TOTAL_RSM2) * 100
        valor_total = valor_rsm2 * AC
        
        # Alternar cor de fundo
        if i % 2 == 1:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(9), Inches(0.35)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_GRAY
            bg.line.fill.background()
        
        values = [nome, format_money_decimal(valor_rsm2), format_percent(perc), format_money(valor_total)]
        for value, x in zip(values, x_positions):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.8), Inches(0.35))
            tf = box.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = NAVY
            tf.paragraphs[0].font.name = 'Calibri'
            if x > 0.5:  # Alinhar números à direita
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        y += 0.4
    
    # Slide 8 - Detalhamento Estrutural
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Detalhamento Estrutural")
    
    # Valores estimados
    supraestrutura_total = 722.66 * AC
    concreto_valor = supraestrutura_total * 0.45
    aco_valor = supraestrutura_total * 0.30
    forma_valor = supraestrutura_total * 0.25
    
    # Estimativas
    concreto_m3 = 3250  # ~0.22 m³/m²
    aco_kg = 292500  # ~90 kg/m³
    forma_m2 = 11850  # ~3.6 m²/m³
    
    content = f"""
CONCRETO (45% da Supraestrutura)
• Volume estimado: {concreto_m3:,.0f} m³
• Índice: 0,22 m³/m² AC
• Custo: {format_money(concreto_valor)}

AÇO (30% da Supraestrutura)
• Peso estimado: {aco_kg:,.0f} kg
• Taxa: ~90 kg/m³ concreto
• Custo: {format_money(aco_valor)}

FORMA (25% da Supraestrutura)
• Área estimada: {forma_m2:,.0f} m²
• Índice: ~3,6 m²/m³ concreto
• Custo: {format_money(forma_valor)}

FUNDAÇÃO (85% da Infraestrutura)
• Tipo: Hélice Contínua
• Estimativa: {format_money(230.40 * AC * 0.85)}
""".strip()
    
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(8.5), Inches(4.2))
    tf = text_box.text_frame
    
    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        if line and not line.startswith('•'):
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.space_before = Pt(8)
        else:
            p.font.size = Pt(14)
            p.font.color.rgb = NAVY
        p.font.name = 'Calibri'
        p.level = 1 if line.startswith('•') else 0
    
    # Slide 9 - Detalhamento Instalações
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Detalhamento de Instalações")
    
    inst_total = 366.96 * AC
    sist_esp_total = 193.39 * AC
    
    # Distribuição instalações
    inst_dist = [
        ("Hidrossanitárias", 33),
        ("Elétricas", 30),
        ("Preventivas", 16),
        ("Gás", 10),
        ("Telecom", 11),
    ]
    
    # Sistemas especiais
    sist_esp = [
        ("Elevadores", 40),
        ("Gerador", 20),
        ("Automação/CFTV", 20),
        ("Pressuriz./Outros", 20),
    ]
    
    # Instalações
    y = 1.2
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(4), Inches(0.3))
    tf = title_box.text_frame
    tf.text = "INSTALAÇÕES"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.name = 'Calibri'
    
    y += 0.4
    for nome, perc in inst_dist:
        valor = inst_total * (perc / 100)
        
        # Barra
        bar_width = (perc / 35) * 3
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(y),
            Inches(bar_width), Inches(0.25)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        
        # Texto
        label_box = slide.shapes.add_textbox(Inches(0.8 + bar_width + 0.1), Inches(y), Inches(3), Inches(0.25))
        tf = label_box.text_frame
        tf.text = f"{nome}: {perc}% ({format_money(valor)})"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        y += 0.35
    
    # Sistemas Especiais
    y = 1.2
    title_box = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(4), Inches(0.3))
    tf = title_box.text_frame
    tf.text = "SISTEMAS ESPECIAIS"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.name = 'Calibri'
    
    y += 0.4
    for nome, perc in sist_esp:
        valor = sist_esp_total * (perc / 100)
        
        # Barra
        bar_width = (perc / 45) * 3
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(y),
            Inches(bar_width), Inches(0.25)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ORANGE
        bar.line.fill.background()
        
        # Texto
        label_box = slide.shapes.add_textbox(Inches(5.5 + bar_width + 0.1), Inches(y), Inches(3.5), Inches(0.25))
        tf = label_box.text_frame
        tf.text = f"{nome}: {perc}% ({format_money(valor)})"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        y += 0.35
    
    # Slide 10 - Detalhamento Acabamentos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Detalhamento de Acabamentos")
    
    acabamentos = [
        ("Rev. Int. Parede", 171.85),
        ("Pisos", 193.69),
        ("Teto", 68.65),
        ("Pintura", 133.12),
        ("Esquadrias", 367.23),
        ("Fachada", 170.21),
    ]
    
    # Header
    headers = ["Item", "R$/m²", "Valor Total"]
    x_positions = [1, 5, 7.5]
    
    y = 1.2
    for header, x in zip(headers, x_positions):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(1.5), Inches(0.3))
        tf = box.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(2), Inches(0.3)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
    
    y = 1.6
    for i, (nome, valor_rsm2) in enumerate(acabamentos):
        valor_total = valor_rsm2 * AC
        
        if i % 2 == 1:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(y),
                Inches(8.5), Inches(0.35)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_GRAY
            bg.line.fill.background()
        
        values = [nome, format_money_decimal(valor_rsm2), format_money(valor_total)]
        for value, x in zip(values, x_positions):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2), Inches(0.35))
            tf = box.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = NAVY
            tf.paragraphs[0].font.name = 'Calibri'
            if x > 1:
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        y += 0.4
    
    # Nota
    note_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.5))
    tf = note_box.text_frame
    tf.text = "Valores ajustam automaticamente com padrão de acabamento"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 11 - Benchmark
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Benchmark — Projetos de Referência")
    
    # Header
    headers = ["Projeto", "Cidade", "AC (m²)", "Pav", "R$/m²", "CUB Ratio"]
    x_positions = [0.5, 2.3, 3.8, 5.3, 6.8, 8.3]
    widths = [1.6, 1.3, 1.3, 1.3, 1.3, 1.5]
    
    y = 1.2
    for header, x, w in zip(headers, x_positions, widths):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
        tf = box.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(0.3)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
    
    # Projetos
    y = 1.6
    for i, (nome, cidade, ac, pav, prazo, rsm2, ratio) in enumerate(BENCHMARK):
        values = [nome, cidade, f"{ac:,.0f}".replace(",", "."), str(pav), 
                 format_money_decimal(rsm2), f"{ratio:.2f}"]
        
        if i % 2 == 1:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5), Inches(y),
                Inches(9.3), Inches(0.3)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_GRAY
            bg.line.fill.background()
        
        for value, x, w in zip(values, x_positions, widths):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
            tf = box.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = NAVY
            tf.paragraphs[0].font.name = 'Calibri'
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        y += 0.35
    
    # Este Projeto (destacado)
    y += 0.15
    values = ["Este Projeto", LOCAL.split('/')[0], f"{AC:,.0f}".replace(",", "."), str(PAVIMENTOS),
             format_money_decimal(TOTAL_RSM2), f"{CUB_RATIO:.3f}"]
    
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(y),
        Inches(9.3), Inches(0.35)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = ORANGE
    bg.line.color.rgb = ORANGE
    bg.line.width = Pt(2)
    
    for value, x, w in zip(values, x_positions, widths):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
        tf = box.text_frame
        tf.text = value
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 12 - Análise do Produto
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Análise do Produto")
    
    # Indicadores calculados
    ac_ur = AC / UR
    custo_ur = TOTAL_VALOR / UR
    vagas_ur = VAGAS / UR
    perc_tipo = (PAVIMENTOS_TIPO / PAVIMENTOS) * 100
    elev_ur = ELEVADORES / UR
    
    indicadores = [
        ("AC/UR", f"{ac_ur:.1f} m²/UR", "90-160 m²", "OK" if 90 <= ac_ur <= 160 else "⚠️"),
        ("Custo/UR", format_money(custo_ur), "—", "—"),
        ("R$/m²", format_money_decimal(TOTAL_RSM2), "—", "—"),
        ("CUB Ratio", f"{CUB_RATIO:.3f}", "1,00-1,50", "OK" if 1.0 <= CUB_RATIO <= 1.5 else "⚠️"),
        ("Vagas/UR", f"{vagas_ur:.1f}", "1,0-2,0", "OK" if 1.0 <= vagas_ur <= 2.0 else "⚠️"),
        ("% Tipo", f"{perc_tipo:.0f}%", "60-80%", "OK" if 60 <= perc_tipo <= 80 else "⚠️"),
        ("Elevador/UR", f"{elev_ur:.3f}", "0,015-0,030", "OK" if 0.015 <= elev_ur <= 0.030 else "⚠️"),
    ]
    
    # Header
    headers = ["Indicador", "Este Projeto", "Faixa Típica", "Status"]
    x_positions = [0.8, 3.2, 5.5, 7.8]
    widths = [2.2, 2, 2, 1.8]
    
    y = 1.2
    for header, x, w in zip(headers, x_positions, widths):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
        tf = box.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(0.3)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
    
    y = 1.6
    for i, (indic, valor, faixa, status) in enumerate(indicadores):
        if i % 2 == 1:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(y),
                Inches(9), Inches(0.35)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_GRAY
            bg.line.fill.background()
        
        values = [indic, valor, faixa, status]
        for value, x, w in zip(values, x_positions, widths):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
            tf = box.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = NAVY
            tf.paragraphs[0].font.name = 'Calibri'
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        y += 0.4
    
    # Slide 13 - Metodologia
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Metodologia do Orçamento Paramétrico")
    
    content = [
        ("Base de Dados", "18 projetos reais de orçamentos executivos da Cartesian"),
        ("Parâmetros", "25 perguntas de briefing ajustáveis (laje, padrão, fundação, etc.)"),
        ("Atualização", "Indexado pelo CUB — ajuste automático ao longo do tempo"),
        ("Validação", "Calibrado com diferença <3% vs orçamentos executivos reais"),
    ]
    
    y = 1.5
    for i, (titulo, desc) in enumerate(content):
        # Número
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), Inches(y),
            Inches(0.4), Inches(0.4)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = ORANGE
        num_circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(0.4), Inches(0.4))
        tf = num_box.text_frame
        tf.text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(1.6), Inches(y), Inches(7.5), Inches(0.2))
        tf = title_box.text_frame
        tf.text = titulo
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        # Descrição
        desc_box = slide.shapes.add_textbox(Inches(1.6), Inches(y + 0.25), Inches(7.5), Inches(0.15))
        tf = desc_box.text_frame
        tf.text = desc
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        
        y += 0.8
    
    # Nota
    note_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.5))
    tf = note_box.text_frame
    tf.text = "Calibrado com dados reais de orçamentos executivos da Cartesian"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 14 - Próximos Passos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    add_header_bar(slide, "Próximos Passos")
    
    passos = [
        "Validar premissas com equipe de projeto",
        "Ajustar briefing conforme definições",
        "Orçamento executivo (quando aplicável)",
    ]
    
    y = 1.8
    for i, passo in enumerate(passos):
        # Número
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.5), Inches(y),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = BLUE
        num_circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(0.5), Inches(0.5))
        tf = num_box.text_frame
        tf.text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Passo
        step_box = slide.shapes.add_textbox(Inches(2.2), Inches(y), Inches(6.5), Inches(0.5))
        tf = step_box.text_frame
        tf.text = passo
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = NAVY
        tf.paragraphs[0].font.name = 'Calibri'
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        y += 0.8
    
    # CTA
    cta_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(0.6))
    tf = cta_box.text_frame
    tf.text = "Planilha interativa disponível para ajustes em tempo real"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 15 - Contato/Encerramento
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    
    # CARTESIAN ENGENHARIA (grande)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "CARTESIAN ENGENHARIA"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Barra laranja
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(2.2),
        Inches(4), Inches(0.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    
    # Contato
    contact_info = [
        "Leonardo Kock Adriano",
        "leonardo@cartesianengenharia.com",
        "cartesianengenharia.com",
    ]
    
    y = 2.8
    for info in contact_info:
        box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.4))
        tf = box.text_frame
        tf.text = info
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = BLUE
        tf.paragraphs[0].font.name = 'Calibri'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        y += 0.45
    
    # Salvar
    output_path = "rozzo-vd-apresentacao.pptx"
    prs.save(output_path)
    print(f"✅ Apresentação criada: {output_path}")
    print(f"   Total de slides: {len(prs.slides)}")
    print(f"   Tamanho: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.2f}\" (16:9)")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
