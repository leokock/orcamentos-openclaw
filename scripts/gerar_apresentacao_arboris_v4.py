"""
Gera a apresentacao .pptx do Arboris paramétrico v4.
Estrutura inspirada no CTN-ALF-SFL (20 slides).
Comparativos com valor MÉDIO (sem nomes de obras/clientes).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ============ Paleta Cartesian ============
AZUL_PRIM = RGBColor(0x1E, 0x52, 0xF5)  # #1E52F5
AZUL_NAVY = RGBColor(0x0F, 0x1C, 0x4A)  # #0F1C4A
LARANJA = RGBColor(0xFF, 0x57, 0x22)    # #FF5722
CINZA_FUNDO = RGBColor(0xF4, 0xF6, 0xF8)  # fundo claro
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_LINHA = RGBColor(0xE8, 0xEE, 0xF4)
CINZA_TXT = RGBColor(0x45, 0x4E, 0x5E)
PRETO = RGBColor(0x1A, 0x1A, 0x1A)

# Formato widescreen 16:9 (13.33 x 7.5 polegadas)
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_blank_slide():
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def add_rect(slide, x, y, w, h, fill_color, line=False, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = fill_color
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=PRETO, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_footer(slide, is_first=False):
    """Rodape discreto com Cartesian a direita e projeto a esquerda."""
    if not is_first:
        add_text(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
                 'Cartesian', size=11, bold=True, color=AZUL_PRIM, align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 - Capa Cartesian
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
# Logo Cartesian centralizado (texto simulando)
add_text(s, Inches(0), Inches(2.8), W, Inches(1.8),
         'Cartesian', size=96, bold=True, color=AZUL_PRIM, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.5), W, Inches(0.5),
         'Coordenação Técnica e Orçamentária', size=20, color=CINZA_TXT, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2 - Capa do Projeto
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
# Lado esquerdo (50%)
add_text(s, Inches(0.8), Inches(1.2), Inches(6.5), Inches(0.9),
         'Orçamento\nParamétrico', size=44, bold=True, color=AZUL_PRIM, font='Calibri')
# Linha divisória
add_rect(s, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.04), AZUL_PRIM)
add_text(s, Inches(0.8), Inches(3.5), Inches(6.5), Inches(0.9),
         'ARBORIS', size=50, bold=True, color=AZUL_NAVY)
add_text(s, Inches(0.8), Inches(4.5), Inches(6.5), Inches(0.4),
         'ARTHEN EMPREENDIMENTOS', size=18, bold=True, color=CINZA_TXT)
add_text(s, Inches(0.8), Inches(5.1), Inches(6.5), Inches(0.4),
         'Itajaí/SC  |  Data-base: Março/2026  |  CUB/SC: R$ 3.028,45/m²',
         size=13, color=CINZA_TXT)
add_text(s, Inches(0.8), Inches(5.5), Inches(6.5), Inches(0.4),
         'Versão 4 híbrida  |  Abril/2026',
         size=11, color=LARANJA, bold=True)
# Lado direito - placeholder render
add_rect(s, Inches(7.5), Inches(0.8), Inches(5.3), Inches(5.9), AZUL_NAVY)
add_text(s, Inches(7.5), Inches(3.3), Inches(5.3), Inches(1.0),
         '[RENDER 3D DA TORRE]', size=16, color=BRANCO, align=PP_ALIGN.CENTER)
add_footer(s)


# =====================================================================
# SLIDE 3 - Jornada da Reunião
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
         'JORNADA DA REUNIÃO', size=32, bold=True, color=AZUL_PRIM)

topics = [
    ('01', 'Dados e Premissas do Projeto'),
    ('02', 'Visão Geral do Orçamento'),
    ('03', 'Composição por Macrogrupo'),
    ('04', 'Especificações'),
    ('05', 'Comparativo com Média de Mercado'),
    ('06', 'Análise de Custo por Área Privativa'),
]
y = Inches(1.7)
for num, title in topics:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.75), BRANCO)
    add_rect(s, Inches(0.8), y, Inches(0.1), Inches(0.75), AZUL_PRIM)
    add_text(s, Inches(1.1), y, Inches(1.5), Inches(0.75),
             num, size=28, bold=True, color=AZUL_PRIM, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), y, Inches(9.5), Inches(0.75),
             title, size=18, bold=True, color=AZUL_NAVY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.85)
add_footer(s)


# =====================================================================
# SLIDE 4 - Características do Empreendimento
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
         'CARACTERÍSTICAS DO EMPREENDIMENTO', size=32, bold=True, color=AZUL_PRIM)

# Tabela à esquerda (54%)
caracteristicas = [
    ('Área Construída', '12.472,98 m²'),
    ('Unidades Residenciais', '98 un.'),
    ('Pavimentos', '24 (19 tipo)'),
    ('Elevadores', '2'),
    ('Vagas de Garagem', '120'),
    ('Tipologia', '3-4 dormitórios, 2 BWCs/apto'),
    ('Padrão de Acabamento', 'Médio-Alto'),
    ('Prazo de Execução', '36 meses'),
]
y = Inches(1.7)
for i, (k, v) in enumerate(caracteristicas):
    add_rect(s, Inches(0.8), y, Inches(6.5), Inches(0.55), AZUL_NAVY)
    add_text(s, Inches(1.0), y, Inches(4.0), Inches(0.55),
             k, size=14, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE, bold=False)
    add_text(s, Inches(5.0), y, Inches(2.1), Inches(0.55),
             v, size=14, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.65)

# Placeholder render direita
add_rect(s, Inches(7.8), Inches(1.7), Inches(5.0), Inches(5.0), AZUL_NAVY)
add_text(s, Inches(7.8), Inches(4.0), Inches(5.0), Inches(0.5),
         '[RENDER 3D TORRE]', size=14, color=BRANCO, align=PP_ALIGN.CENTER)
add_footer(s)


# =====================================================================
# SLIDE 5 - Visão Geral (4 cards)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
         'VISÃO GERAL DO ORÇAMENTO', size=32, bold=True, color=AZUL_PRIM)

cards = [
    ('R$ 44,3 Mi', 'CUSTO TOTAL', 'Valor global do empreendimento'),
    ('R$ 3.551/m²', 'CUSTO POR M²', 'Área construída: 12.473 m²'),
    ('R$ 452 mil', 'CUSTO POR UNIDADE', 'Total de 98 unidades'),
    ('1,17 CUB', 'FATOR CUB/SC', 'CUB Mar/2026: R$ 3.028,45'),
]
xs = [Inches(0.8), Inches(7.0)]
ys = [Inches(1.6), Inches(4.3)]
for i, (big, label, desc) in enumerate(cards):
    x = xs[i % 2]
    y = ys[i // 2]
    add_rect(s, x, y, Inches(5.5), Inches(2.5), AZUL_NAVY)
    add_text(s, x, y+Inches(0.25), Inches(5.5), Inches(0.85),
             big, size=36, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
    add_text(s, x, y+Inches(1.25), Inches(5.5), Inches(0.4),
             label, size=14, bold=True, color=AZUL_PRIM, align=PP_ALIGN.CENTER)
    add_text(s, x, y+Inches(1.75), Inches(5.5), Inches(0.5),
             desc, size=13, color=BRANCO, align=PP_ALIGN.CENTER)
add_footer(s)


# =====================================================================
# SLIDE 6 - Detalhamento do Custo
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'DETALHAMENTO DO CUSTO', size=28, bold=True, color=AZUL_PRIM)

# Tabela
macros = [
    ('GERENCIAMENTO TÉCNICO E ADMINISTRATIVO', 5078256, 11.5, 407),
    ('MOVIMENTAÇÃO DE TERRA', 399135, 0.9, 32),
    ('INFRAESTRUTURA', 2384313, 5.4, 191),
    ('SUPRAESTRUTURA', 9074861, 20.5, 728),
    ('ALVENARIA', 2264768, 5.1, 182),
    ('INSTALAÇÕES ELÉTRICAS, HIDRÁULICAS, GLP E PREVENTIVAS', 4918221, 11.1, 394),
    ('EQUIPAMENTOS E SISTEMAS ESPECIAIS', 1490000, 3.4, 119),
    ('CLIMATIZAÇÃO', 315930, 0.7, 25),
    ('IMPERMEABILIZAÇÃO', 997599, 2.3, 80),
    ('REVESTIMENTOS INTERNOS DE PAREDE', 1601711, 3.6, 128),
    ('REVESTIMENTOS E ACABAMENTOS EM TETO', 905729, 2.0, 73),
    ('PISOS E PAVIMENTAÇÕES', 2380641, 5.4, 191),
    ('PINTURA INTERNA', 1928859, 4.4, 155),
    ('ESQUADRIAS, VIDROS E FERRAGENS', 4053719, 9.2, 325),
    ('LOUÇAS E METAIS', 395136, 0.9, 32),
    ('FACHADA', 3924599, 8.9, 315),
    ('SERVIÇOS COMPLEMENTARES', 1518921, 3.4, 122),
    ('IMPREVISTOS E CONTINGÊNCIAS', 654486, 1.5, 52),
]
total_v = sum(x[1] for x in macros)

# Cabeçalho
hdr_y = Inches(1.2)
add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.4), AZUL_PRIM)
add_text(s, Inches(0.6), hdr_y, Inches(6.6), Inches(0.4),
         'ETAPA', size=11, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.2), hdr_y, Inches(2.2), Inches(0.4),
         'VALOR ORÇADO', size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.6), hdr_y, Inches(1.2), Inches(0.4),
         '%', size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.0), hdr_y, Inches(1.7), Inches(0.4),
         'VALOR / m²', size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

y = Inches(1.6)
for i, (name, val, pct, rsm2) in enumerate(macros):
    bg = BRANCO if i % 2 == 0 else CINZA_LINHA
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.27), bg)
    add_text(s, Inches(0.6), y, Inches(6.6), Inches(0.27),
             name, size=9, color=PRETO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.2), y, Inches(2.2), Inches(0.27),
             f'R$ {val:,.2f}'.replace(',', '.'), size=9, color=PRETO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.6), y, Inches(1.2), Inches(0.27),
             f'{pct:.1f}%', size=9, color=PRETO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(11.0), y, Inches(1.7), Inches(0.27),
             f'R$ {rsm2:.2f}'.replace(',', '.'), size=9, color=PRETO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.27)

# Linha TOTAL
add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.35), AZUL_PRIM)
add_text(s, Inches(0.6), y, Inches(6.6), Inches(0.35),
         'TOTAL', size=12, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.2), y, Inches(2.2), Inches(0.35),
         f'R$ {total_v:,.2f}'.replace(',', '.'), size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.6), y, Inches(1.2), Inches(0.35),
         '100%', size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.0), y, Inches(1.7), Inches(0.35),
         'R$ 3.551,00', size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s)


# =====================================================================
# SLIDE 7 - Composição (top 10 barras + TOP 3 cards)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.8),
         'COMPOSIÇÃO DO ORÇAMENTO', size=30, bold=True, color=AZUL_PRIM)

# Top 10 sorted desc
top10 = sorted([(m[0].split(',')[0].strip(), m[2]) for m in macros], key=lambda x: -x[1])[:10]
# Shorten names
shortnames = {
    'GERENCIAMENTO TÉCNICO E ADMINISTRATIVO': 'Gerenciamento',
    'SUPRAESTRUTURA': 'Supraestrutura',
    'INFRAESTRUTURA': 'Infraestrutura',
    'INSTALAÇÕES ELÉTRICAS': 'Instalações',
    'ESQUADRIAS': 'Esquadrias',
    'FACHADA': 'Fachada',
    'ALVENARIA': 'Alvenaria',
    'SERVIÇOS COMPLEMENTARES': 'Complementares',
    'PISOS E PAVIMENTAÇÕES': 'Pisos',
    'PINTURA INTERNA': 'Pintura',
    'REVESTIMENTOS INTERNOS DE PAREDE': 'Rev. Parede',
    'EQUIPAMENTOS E SISTEMAS ESPECIAIS': 'Sist. Especiais',
    'REVESTIMENTOS E ACABAMENTOS EM TETO': 'Rev. Teto',
    'IMPERMEABILIZAÇÃO': 'Impermeabilização',
    'LOUÇAS E METAIS': 'Louças e Metais',
    'MOVIMENTAÇÃO DE TERRA': 'Mov. Terra',
    'IMPREVISTOS E CONTINGÊNCIAS': 'Imprevistos',
    'CLIMATIZAÇÃO': 'Climatização',
}
top10_named = [(shortnames.get(k, k), v) for k, v in top10]
# Ascending for chart display (maior em baixo)
top10_named = list(reversed(top10_named))
max_pct = max(v for _, v in top10_named)
chart_x = Inches(0.8)
chart_y = Inches(1.4)
chart_w = Inches(7.3)
chart_h = Inches(5.7)
bar_h = (chart_h - Inches(0.3)) / len(top10_named)
for i, (name, v) in enumerate(top10_named):
    y = chart_y + bar_h * i
    # Label
    add_text(s, chart_x, y, Inches(1.8), bar_h,
             name, size=11, color=CINZA_TXT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # Bar
    bar_w = Emu(int((chart_w - Inches(2.0)) * v / max_pct))
    add_rect(s, chart_x + Inches(1.9), y + Emu(int(int(bar_h) * 0.2)),
             bar_w, Emu(int(int(bar_h) * 0.6)), AZUL_PRIM)
    # Value
    add_text(s, chart_x + Inches(1.9) + bar_w + Inches(0.05), y,
             Inches(1.0), bar_h,
             f'{v:.1f}%', size=11, bold=True, color=AZUL_NAVY, anchor=MSO_ANCHOR.MIDDLE)

# TOP 3 cards à direita
top3_data = sorted(macros, key=lambda x: -x[2])[:3]
card_x = Inches(8.4)
card_y = Inches(1.4)
card_w = Inches(4.6)
card_h = Inches(1.75)
reasons = {
    'SUPRAESTRUTURA': 'Laje convencional, torre de 24 pavimentos',
    'GERENCIAMENTO TÉCNICO E ADMINISTRATIVO': '36 meses de obra — equipe técnica completa',
    'INSTALAÇÕES ELÉTRICAS, HIDRÁULICAS, GLP E PREVENTIVAS': 'Elétrica, hidro, PPCI, gás e telecom — 3-4 dorm.',
}
add_text(s, card_x, card_y - Inches(0.1), card_w, Inches(0.4),
         'TOP 3 Macrogrupos', size=16, bold=True, color=AZUL_NAVY)
for i, (name, val, pct, rsm2) in enumerate(top3_data):
    y = card_y + Inches(0.35) + card_h * i
    add_rect(s, card_x, y, card_w, card_h - Inches(0.15), BRANCO)
    add_rect(s, card_x, y, Inches(0.08), card_h - Inches(0.15), LARANJA)
    add_text(s, card_x + Inches(0.25), y + Inches(0.1), Inches(1.8), Inches(0.6),
             f'{pct:.1f}%', size=24, bold=True, color=AZUL_PRIM)
    short = shortnames.get(name, name)
    add_text(s, card_x + Inches(2.0), y + Inches(0.15), Inches(2.5), Inches(0.4),
             short, size=13, bold=True, color=AZUL_NAVY)
    # Value
    add_text(s, card_x + Inches(0.25), y + Inches(0.75), Inches(1.8), Inches(0.35),
             f'R$ {val/1e6:.1f} Mi', size=14, bold=True, color=LARANJA)
    add_text(s, card_x + Inches(2.0), y + Inches(0.75), Inches(2.5), Inches(0.75),
             reasons.get(name, '—'), size=10, color=CINZA_TXT)
add_footer(s)


# =====================================================================
# SLIDES 8-11 - Especificações (4 slides agrupados por magnitude)
# =====================================================================
specs = [
    # slide 8
    [
        ('SUPRAESTRUTURA', '20,5%', 'R$ 9.074.861  |  R$ 728/m²',
         ['Laje convencional maciça (fck30) — consumo 0,25 m³/m²',
          'Aço CA-50 em 106 kg/m² (compatível com pé-direito 3,00m)',
          'Inclui: concreto, formas, armação, escoramento e MO empreitada']),
        ('GERENCIAMENTO TÉCNICO', '11,5%', 'R$ 5.078.256  |  R$ 407/m²',
         ['Prazo de 36 meses — proporcional ao custo de gerenciamento',
          'Inclui: projetos, taxas, ensaios, administração, equipe técnica, segurança, canteiro',
          'Cremalheira própria e minigrua em regime de locação']),
        ('INSTALAÇÕES', '11,1%', 'R$ 4.918.221  |  R$ 394/m²',
         ['Inclui: instalações hidrossanitárias, elétricas, preventivas, gás e telecom',
          'Louças e metais com entrega completa em área comum, bacia sanitária e acabamento de registro em áreas privativas']),
        ('ESQUADRIAS, VIDROS E FERRAGENS', '9,2%', 'R$ 4.053.719  |  R$ 325/m²',
         ['Esquadrias com pintura eletrostática',
          'Inclui: esquadrias de alumínio, vidros, ferragens e instalação',
          'Inclui: esquadrias de ferro, madeira, corrimãos e serralherias']),
    ],
    # slide 9
    [
        ('FACHADA', '8,9%', 'R$ 3.924.599  |  R$ 315/m²',
         ['Revestimento cerâmico de fachada — grande formato',
          'Inclui: material + MO empreitada + balancim fachadeiro',
          'Índice de fachada 1,55 m²/m² de área construída (torre vertical 24 pav.)']),
        ('INFRAESTRUTURA', '5,4%', 'R$ 2.384.313  |  R$ 191/m²',
         ['Fundação em hélice contínua Ø40cm — 190 estacas, 20m médio',
          'Blocos, baldrames e contenção de subsolo (1 nível)',
          'Inclui: perfuração, concreto, aço, forma e MO empreitada']),
        ('PISOS E PAVIMENTAÇÕES', '5,4%', 'R$ 2.380.641  |  R$ 191/m²',
         ['Contrapiso autonivelante + porcelanato principal em áreas privativas',
          'Cimentado liso em garagens e granito em halls comuns',
          'Inclui: rodapés, soleiras e MO de assentamento']),
        ('ALVENARIA', '5,1%', 'R$ 2.264.768  |  R$ 182/m²',
         ['Bloco cerâmico 14cm em áreas secas',
          'Drywall ST em áreas internas e RU em áreas molhadas',
          '82% da vedação em drywall (isolamento acústico Médio-Alto)']),
    ],
    # slide 10
    [
        ('PINTURA INTERNA', '4,4%', 'R$ 1.928.859  |  R$ 155/m²',
         ['Pintura completa em paredes (selador + massa + tinta)',
          'Tinta acrílica em 3 demãos em paredes e 2 demãos em tetos',
          'Textura acrílica em garagens, escadas e áreas técnicas']),
        ('REV. INTERNO DE PAREDE', '3,6%', 'R$ 1.601.711  |  R$ 128/m²',
         ['Argamassado em massa única e estucamento em estruturas',
          'Porcelanato padrão em BWCs, cozinhas e áreas comuns',
          'Granito em bancadas de BWC e cozinha']),
        ('COMPLEMENTARES', '3,4%', 'R$ 1.518.921  |  R$ 122/m²',
         ['Inclui: mobiliário e decoração, paisagismo, comunicação visual',
          'Limpeza final, ligações definitivas e desmobilização',
          'Cobertura, pavimentação externa e equipamentos de lazer']),
        ('EQUIPAMENTOS E SIST. ESPECIAIS', '3,4%', 'R$ 1.490.000  |  R$ 119/m²',
         ['2 elevadores (1 social + 1 serviço) com capacidade para 19 pav. tipo',
          'Gerador dedicado + automação (BMS + sensores) + CFTV + interfonia',
          'Piscina, SPDA, bombas de recalque e quadros de comando']),
    ],
    # slide 11
    [
        ('IMPERMEABILIZAÇÃO', '2,3%', 'R$ 997.599  |  R$ 80/m²',
         ['Regularização + manta asfáltica 4mm em áreas externas, terraços e piscina',
          'Argamassa polimérica em BWCs, lavabos e áreas técnicas',
          'Subsolo (piso enterrado) e cobertura/reservatórios']),
        ('REV. E ACABAMENTOS EM TETO', '2,0%', 'R$ 905.729  |  R$ 73/m²',
         ['Estucamento em garagens e escadas',
          'Forro de gesso acartonado nas áreas privativas e comuns',
          'Forro RU em BWCs e áreas molhadas']),
        ('IMPREVISTOS E CONTINGÊNCIAS', '1,5%', 'R$ 654.486  |  R$ 52/m²',
         ['Percentual aplicado sobre o subtotal do orçamento',
          'Cobre desvios menores entre paramétrico e executivo',
          'Padrão Cartesian para orçamentos pré-executivos']),
        ('MOV. TERRA + CLIMATIZAÇÃO', '1,6%', 'R$ 715.065  |  R$ 57/m²',
         ['Mov. Terra: escavação mecanizada com bota-fora 30% — 1 subsolo',
          'Climatização: entrega com infraestrutura (drenos + eletrodutos)',
          'Exaustão mecânica em BWCs enclausurados e churrasqueiras']),
    ],
]

for block in specs:
    s = add_blank_slide()
    add_rect(s, 0, 0, W, H, CINZA_FUNDO)
    add_text(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.8),
             'ESPECIFICAÇÕES', size=32, bold=True, color=AZUL_PRIM)
    y = Inches(1.5)
    for title, pct, val, bullets in block:
        add_rect(s, Inches(0.8), y, Inches(0.1), Inches(1.35), LARANJA)
        add_rect(s, Inches(0.9), y, Inches(11.7), Inches(1.35), BRANCO)
        add_text(s, Inches(1.0), y+Inches(0.05), Inches(6.5), Inches(0.4),
                 title, size=14, bold=True, color=AZUL_NAVY)
        add_text(s, Inches(7.3), y+Inches(0.05), Inches(1.5), Inches(0.4),
                 pct, size=15, bold=True, color=AZUL_PRIM, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(8.9), y+Inches(0.05), Inches(3.6), Inches(0.4),
                 val, size=11, color=CINZA_TXT, align=PP_ALIGN.RIGHT)
        # bullets
        by = y + Inches(0.5)
        for b in bullets:
            add_text(s, Inches(1.0), by, Inches(11.5), Inches(0.3),
                     '• ' + b, size=10, color=PRETO)
            by += Inches(0.28)
        y += Inches(1.45)
    add_footer(s)


# =====================================================================
# SLIDE 12 - Comparativo com Média de Mercado (SEM nomes de obras)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'COMPARATIVO COM MÉDIA DE MERCADO', size=28, bold=True, color=AZUL_PRIM)
add_text(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
         'Referência: obras de padrão Médio-Alto e Alto da base Cartesian (amostra de 8 empreendimentos)',
         size=11, color=CINZA_TXT)

# Tabela de comparação — sem nomes
comp = [
    ('Mov. Terra', 32, '12 – 31', 20, '+6%'),
    ('Infraestrutura', 191, '171 – 237', 196, '-3%'),
    ('Supraestrutura', 728, '650 – 736', 710, '+3%'),
    ('Alvenaria', 182, '155 – 197', 168, '+8%'),
    ('Instalações', 394, '310 – 445', 330, '+19%'),
    ('Sist. Especiais', 119, '111 – 241', 173, '-31%'),
    ('Climatização', 25, '28 – 86', 42, '-40%'),
    ('Impermeabilização', 80, '35 – 64', 57, '+40%'),
    ('Rev. Parede', 128, '196 – 496', 238, '-46%'),
    ('Teto', 73, '47 – 88', 61, '+20%'),
    ('Pisos', 191, '134 – 327', 284, '-33%'),
    ('Pintura', 155, '137 – 178', 145, '+7%'),
    ('Esquadrias', 325, '265 – 312', 279, '+16%'),
    ('Louças e Metais', 32, '24 – 28', 27, '+19%'),
    ('Fachada', 315, '110 – 259', 149, '+111%'),
    ('Complementares', 122, '104 – 197', 145, '-16%'),
    ('Imprevistos', 52, '42 – 120', 48, '+8%'),
    ('TOTAL (s/ ger.)', 3144, '2.600 – 3.450', 3030, '+4%'),
]
hdr_y = Inches(1.6)
add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.4), AZUL_PRIM)
heads = [('ETAPA', 0.6, 3.5, PP_ALIGN.LEFT),
         ('ARBORIS v4  R$/m²', 4.1, 2.5, PP_ALIGN.CENTER),
         ('FAIXA REFERÊNCIA', 6.7, 2.5, PP_ALIGN.CENTER),
         ('MÉDIA REF.', 9.3, 1.7, PP_ALIGN.CENTER),
         ('VARIAÇÃO', 11.0, 1.8, PP_ALIGN.CENTER)]
for txt, x, w, al in heads:
    add_text(s, Inches(x), hdr_y, Inches(w), Inches(0.4),
             txt, size=11, bold=True, color=BRANCO, align=al, anchor=MSO_ANCHOR.MIDDLE)

y = Inches(2.0)
for i, (name, arb, faixa, med, var) in enumerate(comp):
    is_total = name.startswith('TOTAL')
    bg = AZUL_PRIM if is_total else (BRANCO if i % 2 == 0 else CINZA_LINHA)
    txt_color = BRANCO if is_total else PRETO
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.27), bg)
    add_text(s, Inches(0.6), y, Inches(3.5), Inches(0.27),
             name, size=10, color=txt_color, bold=is_total, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.1), y, Inches(2.5), Inches(0.27),
             f'{arb:,}'.replace(',', '.'), size=10, color=txt_color, bold=is_total, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.7), y, Inches(2.5), Inches(0.27),
             faixa, size=10, color=txt_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.3), y, Inches(1.7), Inches(0.27),
             f'{med:,}'.replace(',', '.'), size=10, color=txt_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Variação com cor
    var_color = LARANJA if is_total else (AZUL_PRIM if var.startswith('-') else PRETO)
    if is_total:
        var_color = BRANCO
    add_text(s, Inches(11.0), y, Inches(1.8), Inches(0.27),
             var, size=10, bold=True, color=var_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.27)
add_footer(s)


# =====================================================================
# (Slide removido: Comparativo com Versão Anterior — interno, não vai pro cliente)
# =====================================================================


# =====================================================================
# SLIDE 13 - Comparação com Média de Mercado (padrão Médio-Alto)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'COMPARAÇÃO COM MÉDIA DE MERCADO', size=28, bold=True, color=AZUL_PRIM)
add_text(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
         'Médias por padrão — base Cartesian (60 obras Médio-Alto + 57 Alto)',
         size=12, color=CINZA_TXT)

# Tabela 4 cols: Etapa | Arboris v4 | Média M-A | Média Alto
tabela_comp = [
    ('Gerenciamento', 407, 475, 480),
    ('Mov. Terra', 32, 115, 100),
    ('Infraestrutura', 191, 196, 210),
    ('Supraestrutura', 728, 655, 780),
    ('Alvenaria', 182, 140, 180),
    ('Instalações', 394, 332, 380),
    ('Sist. Especiais', 119, 173, 230),
    ('Climatização', 25, 31, 60),
    ('Impermeabilização', 80, 60, 70),
    ('Rev. Parede', 128, 61, 180),
    ('Teto', 73, 60, 75),
    ('Pisos', 191, 199, 280),
    ('Pintura', 155, 125, 160),
    ('Esquadrias', 325, 301, 360),
    ('Louças e Metais', 32, 60, 90),
    ('Fachada', 315, 117, 220),
    ('Complementares', 122, 201, 250),
    ('Imprevistos', 52, 50, 55),
    ('TOTAL', 3551, 3349, 4156),
]

hdr_y = Inches(1.7)
add_rect(s, Inches(2.0), hdr_y, Inches(9.3), Inches(0.4), AZUL_NAVY)
add_text(s, Inches(2.1), hdr_y, Inches(3.5), Inches(0.4),
         'ETAPA', size=12, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.6), hdr_y, Inches(2.0), Inches(0.4),
         'ARBORIS v4', size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.6), hdr_y, Inches(2.0), Inches(0.4),
         'MÉDIA MÉDIO-ALTO', size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.6), hdr_y, Inches(1.7), Inches(0.4),
         'MÉDIA ALTO', size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

y = Inches(2.1)
for i, (name, arb, ma, a) in enumerate(tabela_comp):
    is_total = name == 'TOTAL'
    bg = AZUL_PRIM if is_total else (BRANCO if i % 2 == 0 else CINZA_LINHA)
    col_txt = BRANCO if is_total else PRETO
    add_rect(s, Inches(2.0), y, Inches(9.3), Inches(0.23), bg)
    add_text(s, Inches(2.1), y, Inches(3.5), Inches(0.23),
             name, size=10, color=col_txt, bold=is_total, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.6), y, Inches(2.0), Inches(0.23),
             f'{arb:,}'.replace(',', '.'), size=10, color=AZUL_PRIM if not is_total else BRANCO, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.6), y, Inches(2.0), Inches(0.23),
             f'{ma:,}'.replace(',', '.'), size=10, color=col_txt, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.6), y, Inches(1.7), Inches(0.23),
             f'{a:,}'.replace(',', '.'), size=10, color=col_txt, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.23)

add_text(s, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
         'Valores em R$/m² de área construída. Médias indexadas ao CUB/SC Mar/2026.',
         size=10, color=CINZA_TXT)
add_footer(s)


# =====================================================================
# SLIDE 15 - Custo por Área Privativa (comparativo sem nomes)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'ANÁLISE DE CUSTO POR ÁREA PRIVATIVA', size=28, bold=True, color=AZUL_PRIM)

# Tabela
priv_data = [
    ('ARBORIS (este)', '~7.500 m²  (estimado)', 'R$ 5.900', True),
    ('Média Médio-Alto (n=37)', '— ', 'R$ 5.650', False),
    ('Média Alto (n=23)', '—', 'R$ 7.100', False),
]
tab_y = Inches(1.7)
add_rect(s, Inches(1.5), tab_y, Inches(10.3), Inches(0.45), AZUL_NAVY)
add_text(s, Inches(1.7), tab_y, Inches(3.8), Inches(0.45),
         'EMPREENDIMENTO', size=13, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.5), tab_y, Inches(3.5), Inches(0.45),
         'ÁREA PRIVATIVA', size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.0), tab_y, Inches(2.8), Inches(0.45),
         'R$/m² PRIV.', size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

y = tab_y + Inches(0.5)
for name, ap, rsm2, highlight in priv_data:
    bg = CINZA_LINHA if highlight else BRANCO
    txt_color = AZUL_PRIM if highlight else PRETO
    add_rect(s, Inches(1.5), y, Inches(10.3), Inches(0.5), bg)
    add_text(s, Inches(1.7), y, Inches(3.8), Inches(0.5),
             name, size=13, bold=highlight, color=txt_color, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.5), y, Inches(3.5), Inches(0.5),
             ap, size=12, color=PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.0), y, Inches(2.8), Inches(0.5),
             rsm2, size=14, bold=True, color=AZUL_PRIM if highlight else PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)

# Gráfico de barras simples
chart_y = Inches(4.5)
vals = [5900, 5650, 7100]
labels = ['ARBORIS', 'Média M-A', 'Média Alto']
maxv = max(vals)
chart_w = Inches(8)
chart_x = Inches(2.7)
bar_w = Inches(2.2)
for i, (v, lbl) in enumerate(zip(vals, labels)):
    x = chart_x + Inches(i*2.6)
    bh = Inches(2.0 * v / maxv)
    y_bar = chart_y + Inches(2.0) - bh
    add_rect(s, x, y_bar, bar_w, bh, AZUL_PRIM if i == 0 else CINZA_TXT)
    add_text(s, x, y_bar - Inches(0.3), bar_w, Inches(0.3),
             f'R$ {v:,}'.replace(',', '.'), size=12, bold=True, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chart_y + Inches(2.1), bar_w, Inches(0.3),
             lbl, size=11, bold=(i==0), color=PRETO, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         'Valores de referência. A área privativa exata do Arboris depende do levantamento arquitetônico final.',
         size=9, color=CINZA_TXT)
add_footer(s)


# =====================================================================
# SLIDE 16 - Custo por Tipologia
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'CUSTO POR TIPOLOGIA DE UNIDADE', size=28, bold=True, color=AZUL_PRIM)

# Dados indicativos (com nota de que depende do levantamento)
tipos = [
    ('Tipologia A (3 dorm.)', '75,0 m²', 'R$ 442.500'),
    ('Tipologia B (3 dorm. c/ suíte)', '82,0 m²', 'R$ 483.800'),
    ('Tipologia C (4 dorm.)', '92,0 m²', 'R$ 542.800'),
    ('Studio/Apartamento compacto', '55,0 m²', 'R$ 324.500'),
]

tab_y = Inches(1.7)
add_rect(s, Inches(1.5), tab_y, Inches(10.3), Inches(0.45), AZUL_NAVY)
for txt, x, w in [('TIPOLOGIA', 1.7, 4.0), ('ÁREA PRIVATIVA', 6.0, 3.0), ('CUSTO ESTIMADO', 9.0, 2.8)]:
    add_text(s, Inches(x), tab_y, Inches(w), Inches(0.45),
             txt, size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

y = tab_y + Inches(0.5)
for i, (tipo, ap, custo) in enumerate(tipos):
    bg = BRANCO if i % 2 == 0 else CINZA_LINHA
    add_rect(s, Inches(1.5), y, Inches(10.3), Inches(0.45), bg)
    add_text(s, Inches(1.7), y, Inches(4.0), Inches(0.45),
             tipo, size=12, bold=True, color=PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.0), y, Inches(3.0), Inches(0.45),
             ap, size=12, color=PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.0), y, Inches(2.8), Inches(0.45),
             custo, size=12, color=PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.5)

# Resumo
add_rect(s, Inches(1.5), y + Inches(0.3), Inches(10.3), Inches(0.6), AZUL_PRIM)
add_text(s, Inches(1.7), y + Inches(0.3), Inches(10.1), Inches(0.6),
         f'Custo médio por unidade: R$ 451.907  |  R$/m² privativo: R$ 5.900  (baseado em área privativa ~7.500 m²)',
         size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
         'Valores indicativos. Tipologias e áreas privativas serão confirmadas após levantamento arquitetônico final.',
         size=9, color=CINZA_TXT)
add_footer(s)


# =====================================================================
# SLIDE 17 - Planejamento Macro (Gantt 36 meses)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'PLANEJAMENTO MACRO', size=28, bold=True, color=AZUL_PRIM)

# 36 meses
N_MONTHS = 36
gantt = [
    ('Mov. Terra',          [1, 3]),
    ('Infraestrutura',      [2, 7]),
    ('Supraestrutura',      [5, 17]),
    ('Alvenaria',           [8, 22]),
    ('Instalações',         [9, 25]),
    ('Sist. Especiais',     [14, 26]),
    ('Climatização',        [14, 28]),
    ('Impermeabilização',   [12, 27]),
    ('Rev. Piso+Parede',    [14, 30]),
    ('Rev. Teto',           [15, 28]),
    ('Pintura',             [16, 32]),
    ('Esquadrias',          [18, 30]),
    ('Fachada',             [14, 28]),
    ('Complementares',      [28, 36]),
]

table_x = Inches(0.5)
table_y = Inches(1.3)
table_w = Inches(12.3)
label_w = Inches(2.5)
months_w = table_w - label_w
month_w = Emu(int(int(months_w) / N_MONTHS))
row_h = Inches(0.3)

# Header meses
for m in range(N_MONTHS):
    x = table_x + label_w + month_w * m
    add_rect(s, x, table_y, month_w, row_h, AZUL_PRIM)
    add_text(s, x, table_y, month_w, row_h,
             f'M{m+1}', size=7, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Linhas
for i, (name, (start, end)) in enumerate(gantt):
    y = table_y + row_h + row_h * i
    add_rect(s, table_x, y, label_w, row_h, CINZA_LINHA if i % 2 == 0 else BRANCO)
    add_text(s, table_x + Inches(0.05), y, label_w, row_h,
             name, size=8, color=PRETO, anchor=MSO_ANCHOR.MIDDLE)
    for m in range(N_MONTHS):
        x = table_x + label_w + month_w * m
        if start <= m+1 <= end:
            add_rect(s, x, y, month_w, row_h, AZUL_PRIM)
        else:
            add_rect(s, x, y, month_w, row_h, BRANCO)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         'Cronograma macro indicativo — 36 meses de execução',
         size=9, color=CINZA_TXT)
add_footer(s)


# =====================================================================
# SLIDE 18 - Curva de Distribuição de Custos (% por mês + R$ por mês)
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'CURVA DE DISTRIBUIÇÃO DE CUSTOS', size=28, bold=True, color=AZUL_PRIM)

# Distribuição aproximada dos custos ao longo dos 36 meses (curva S)
# Pesos relativos
import math
TOTAL_V4 = 44286884
# Curva S: início baixo, pico M16-20, cauda longa
pesos = []
for m in range(1, N_MONTHS+1):
    # Curva normal aproximada centrada em M17
    p = math.exp(-((m - 17)**2) / (2 * 7**2))
    pesos.append(p)
soma = sum(pesos)
pct = [p / soma * 100 for p in pesos]
valores = [p / 100 * TOTAL_V4 for p in pct]

# Gráfico % barras no topo
chart1_y = Inches(1.4)
chart1_h = Inches(2.7)
chart1_x = Inches(0.5)
chart1_w = Inches(12.3)
maxp = max(pct)
bar_w2 = Emu(int(int(chart1_w - Inches(0.5)) / N_MONTHS))
for m in range(N_MONTHS):
    v = pct[m]
    bh = Inches(2.3 * v / maxp)
    x = chart1_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart1_y + Inches(2.4) - bh, bar_w2 - Emu(50000), bh, AZUL_PRIM)
    add_text(s, x - Emu(30000), chart1_y + Inches(2.4) - bh - Inches(0.25), bar_w2 + Emu(60000), Inches(0.2),
             f'{v:.1f}%', size=6, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chart1_y + Inches(2.45), bar_w2, Inches(0.15),
             f'M{m+1}', size=6, color=CINZA_TXT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3),
         'Distribuição mensal (%)', size=12, bold=True, color=AZUL_NAVY)

# Gráfico R$ barras no meio
chart2_y = Inches(4.4)
maxv = max(valores)
for m in range(N_MONTHS):
    v = valores[m]
    bh = Inches(2.0 * v / maxv)
    x = chart1_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart2_y + Inches(2.0) - bh, bar_w2 - Emu(50000), bh, LARANJA)
    add_text(s, x - Emu(30000), chart2_y + Inches(2.0) - bh - Inches(0.22), bar_w2 + Emu(60000), Inches(0.2),
             f'{v/1e6:.1f}', size=6, color=AZUL_NAVY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(4.1), Inches(5), Inches(0.3),
         'Valor mensal (R$ Mi)', size=12, bold=True, color=AZUL_NAVY)
add_footer(s)


# =====================================================================
# SLIDE 19 - Curva acumulada
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'CURVA ACUMULADA DE CUSTOS', size=28, bold=True, color=AZUL_PRIM)

# Acumulados
pct_acum = []
tot = 0
for p in pct:
    tot += p
    pct_acum.append(tot)
val_acum = [p / 100 * TOTAL_V4 for p in pct_acum]

# Plot linha (simples - usar barras altas pra simular)
chart_y = Inches(1.4)
chart_h = Inches(2.7)
chart_x = Inches(0.5)
chart_w = Inches(12.3)
bar_w2 = Emu(int(int(chart_w - Inches(0.5)) / N_MONTHS))
for m in range(N_MONTHS):
    v = pct_acum[m]
    bh = Inches(2.3 * v / 100)
    x = chart_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart_y + Inches(2.4) - bh, bar_w2 - Emu(30000), bh, AZUL_PRIM)
    if m % 3 == 0:
        add_text(s, x - Emu(30000), chart_y + Inches(2.4) - bh - Inches(0.25),
                 bar_w2 + Emu(60000), Inches(0.2),
                 f'{v:.0f}%', size=7, bold=True, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chart_y + Inches(2.45), bar_w2, Inches(0.15),
             f'M{m+1}', size=6, color=CINZA_TXT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3),
         'Acumulado (%)', size=12, bold=True, color=AZUL_NAVY)

# Acumulado R$
chart2_y = Inches(4.4)
for m in range(N_MONTHS):
    v = val_acum[m]
    bh = Inches(2.0 * v / TOTAL_V4)
    x = chart_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart2_y + Inches(2.0) - bh, bar_w2 - Emu(30000), bh, LARANJA)
    if m % 6 == 0 or m == N_MONTHS - 1:
        add_text(s, x - Emu(40000), chart2_y + Inches(2.0) - bh - Inches(0.25),
                 bar_w2 + Emu(80000), Inches(0.2),
                 f'{v/1e6:.0f}', size=7, bold=True, color=AZUL_NAVY, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(4.1), Inches(5), Inches(0.3),
         'Acumulado (R$ Mi)', size=12, bold=True, color=AZUL_NAVY)
add_footer(s)


# =====================================================================
# SLIDES DE JUSTIFICATIVA — Macrogrupos acima da média Médio-Alto
# Divididos em 3 slides (3 itens/slide)
# =====================================================================

# Lista de justificativas — (macrogrupo, R$/m² Arboris, R$/m² média MA, delta %, título, texto)
justificativas = [
    ('SUPRAESTRUTURA', 728, 655, '+11%',
     'Torre vertical de 24 pavimentos com laje convencional maciça',
     [
         'Sistema estrutural: laje maciça em concreto armado convencional (fck30), com consumo de 0,25 m³/m² — no topo da faixa das obras Médio-Alto',
         'Aço CA-50 em 106 kg/m² — compatível com pé-direito padrão de 3,00m e 19 pavimentos tipo',
         'Forma em compensado plastificado com índice de 7,12 m²/m² — padrão para torres verticais',
         'MO estrutural empreitada em R$ 185/m² (mediana da base Cartesian) + escoramento metálico em locação',
         'A escolha da laje maciça (em vez de nervurada) eleva o custo estrutural em ~10% mas facilita o cumprimento de prazos e melhora o desempenho acústico entre apartamentos',
     ]),
    ('ALVENARIA', 182, 140, '+30%',
     '82% das vedações em drywall (isolamento acústico premium)',
     [
         'Especificação com drywall ST em paredes internas (7.016 m²) e drywall RU em áreas molhadas (2.245 m²)',
         'Bloco cerâmico 14cm usado apenas em perímetros externos (11.226 m²)',
         'Drywall ST a R$ 156/m² e RU a R$ 195/m² refletem chapas duplas com lã mineral — padrão Médio-Alto para isolamento acústico de classe A',
         'A maioria das obras da base usa apenas bloco cerâmico (R$ 33/m²), daí a diferença de +30% na média',
         'O investimento em drywall heavy entrega redução significativa de ruído de impacto e aéreo entre unidades — diferencial comercial do empreendimento',
     ]),
    ('INSTALAÇÕES', 394, 332, '+19%',
     'Tipologia 3-4 dormitórios com 2 BWCs/apto e entrega completa',
     [
         'Tipologia 3-4 dormitórios aplica fator de 1,15 sobre quantidade de pontos elétricos vs obras com tipologias menores',
         'Apartamentos com 2 BWCs elevam o fator hidráulico em 1,15 — mais pontos de água/esgoto/ventilação',
         'Entrega completa com louças e metais instalados adiciona MO de instalação final',
         'Escopo inclui elétrica, hidrossanitárias, preventivas, gás canalizado e telecom integrados',
         'Aço inox e tubulação PPR soldada em coluna principal — padrão Médio-Alto',
     ]),
    ('IMPERMEABILIZAÇÃO', 80, 60, '+33%',
     'Cobertura ampla: subsolo + 98 BWCs + piscina + terraços',
     [
         'Índice de impermeabilização em 0,45 m²/m² AC — topo da faixa para obras com 1 subsolo',
         'Piso enterrado do subsolo recebe manta asfáltica 4mm com regularização e proteção mecânica',
         '98 banheiros privativos + áreas comuns com argamassa polimérica em todas as áreas molhadas',
         'Piscina condominial, cobertura, reservatórios superiores e terraços de unidades recebem manta asfáltica',
         'O valor acima da média reflete o escopo ampliado de áreas molhadas — não sobrepreço',
     ]),
    ('REV. INTERNO DE PAREDE', 128, 61, '+110%',
     'Entrega completa com porcelanato, cerâmica e granito (não é shell)',
     [
         'Escopo amplo: reboco + chapisco + cerâmica 30×60 em BWCs/cozinhas + porcelanato em áreas comuns + granito em bancadas',
         'Cerâmica BWC a R$ 48/m² e cozinha a R$ 42/m² — mediana da base Cartesian para padrão Médio-Alto',
         'Porcelanato em áreas comuns a R$ 85/m² — corredor, halls e área social com acabamento premium',
         'Granito em bancadas de BWC e cozinha (343 m linear) — padrão de entrega completa',
         'A média de R$ 61/m² está puxada por obras com entrega shell (sem pisos/revestimentos privativos) — Arboris entrega completa',
     ]),
    ('TETO', 73, 60, '+22%',
     'Forro de gesso acartonado em todas as áreas (comum + privativa)',
     [
         'Forro ST em 10.128 m² (áreas privativas e comuns secas) + Forro RU em 2.170 m² (BWCs e áreas molhadas)',
         'PU de R$ 28/m² para ST e R$ 35/m² para RU — mediana da base',
         'Estucamento em garagens e escadas — acabamento padronizado Médio-Alto',
         'Obras com entrega shell costumam não incluir forro privativo, baixando a média',
     ]),
    ('PINTURA INTERNA', 155, 125, '+24%',
     'Entrega completa: selador + massa PVA + 3 demãos de acrílica',
     [
         'Paredes: selador + massa PVA + 3 demãos de tinta acrílica — padrão Médio-Alto para acabamento liso',
         'Tetos: 2 demãos de acrílica com selador',
         'Texturas acrílicas em garagens, escadas e áreas técnicas',
         'Demarcação epóxi de vagas e pintura de piso em áreas técnicas',
         'MO empreitada representa 66,2% do custo da disciplina — padrão de mão de obra especializada',
     ]),
    ('ESQUADRIAS, VIDROS E FERRAGENS', 325, 301, '+8%',
     'Esquadrias alumínio com pintura eletrostática + serralheria completa',
     [
         'Esquadrias de alumínio com pintura eletrostática — padrão Médio-Alto (não anodizado comum)',
         'Vidros temperados 6mm/8mm em janelas e portas de sacada',
         'Ferragens de abertura e fechamento de marca consolidada',
         'Serralheria interna completa: corrimãos de escada, guarda-corpos, portões de garagem, grades técnicas',
         'Acima da média por 8% — dentro da faixa esperada para o padrão, com margem de 5-10%',
     ]),
    ('FACHADA', 315, 117, '+169%',
     'Cerâmica de fachada em torre vertical de 24 pavimentos',
     [
         'Revestimento em placa cerâmica de fachada — especificação premium (não textura projetada)',
         'PU material R$ 150/m² (cerâmica grande formato) + MO empreitada R$ 35/m² + balancim R$ 18/m²',
         'Índice de fachada em 1,55 m²/m² AC — padrão para torre vertical de 24 pavimentos (obras baixas têm 0,8-1,2)',
         'A maioria das obras da base usa textura projetada (R$ 110-150/m² AC) — Arboris especificou cerâmica, que tem custo estrutural superior mas durabilidade e apelo comercial muito maiores',
         'Obras similares com fachada cerâmica na base Cartesian operam em R$ 259-300/m² AC — Arboris está dentro dessa faixa específica',
     ]),
]

# Divide em 3 slides: 3 itens cada (3+3+3)
n_per_slide = 3
grupos = [justificativas[i:i+n_per_slide] for i in range(0, len(justificativas), n_per_slide)]

for idx, grupo in enumerate(grupos):
    s = add_blank_slide()
    add_rect(s, 0, 0, W, H, CINZA_FUNDO)
    title = f'JUSTIFICATIVAS — ITENS ACIMA DA MÉDIA ({idx+1}/{len(grupos)})'
    add_text(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.6),
             title, size=24, bold=True, color=AZUL_PRIM)
    add_text(s, Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.35),
             'Macrogrupos com valor superior à média de mercado Médio-Alto e as razões técnicas',
             size=11, color=CINZA_TXT)

    # 3 blocos horizontais ou verticais conforme altura
    y = Inches(1.5)
    block_h = Inches(1.85)
    for mg, arb, media, delta, subtitle, bullets in grupo:
        # Fundo do bloco
        add_rect(s, Inches(0.5), y, Inches(12.3), block_h, BRANCO)
        add_rect(s, Inches(0.5), y, Inches(0.12), block_h, LARANJA)
        # Header do bloco: nome + valores
        add_text(s, Inches(0.75), y + Inches(0.05), Inches(6.5), Inches(0.4),
                 mg, size=13, bold=True, color=AZUL_NAVY)
        # Pastilha Arboris vs Média com delta
        add_rect(s, Inches(7.2), y + Inches(0.08), Inches(5.3), Inches(0.38), CINZA_FUNDO)
        add_text(s, Inches(7.3), y + Inches(0.08), Inches(2.5), Inches(0.38),
                 f'Arboris: R$ {arb}/m²', size=10, bold=True, color=AZUL_PRIM, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.8), y + Inches(0.08), Inches(1.6), Inches(0.38),
                 f'Média: R$ {media}/m²', size=10, color=CINZA_TXT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(11.4), y + Inches(0.08), Inches(1.05), Inches(0.38),
                 delta, size=11, bold=True, color=LARANJA, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Subtítulo
        add_text(s, Inches(0.75), y + Inches(0.5), Inches(12), Inches(0.3),
                 subtitle, size=11, bold=True, color=AZUL_PRIM)
        # Bullets (limitar a 3 bullets principais por slide pra caber)
        by = y + Inches(0.85)
        for b in bullets[:3]:
            add_text(s, Inches(0.85), by, Inches(11.8), Inches(0.32),
                     '• ' + b, size=9, color=PRETO)
            by += Inches(0.32)
        y += block_h + Inches(0.1)
    add_footer(s)


# =====================================================================
# SLIDE FINAL - Entregáveis
# =====================================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
         'ENTREGÁVEIS', size=32, bold=True, color=AZUL_PRIM)

entregaveis = [
    ('01', 'Apresentação'),
    ('02', 'Orçamento Analítico (Excel)'),
    ('03', 'Planejamento Macro e Distribuição Físico-Financeira'),
]
y = Inches(2.0)
for num, t in entregaveis:
    add_rect(s, Inches(0.8), y, Inches(7.5), Inches(0.85), BRANCO)
    add_rect(s, Inches(0.8), y, Inches(0.12), Inches(0.85), AZUL_PRIM)
    add_text(s, Inches(1.1), y, Inches(1.5), Inches(0.85),
             num, size=30, bold=True, color=AZUL_PRIM, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), y, Inches(5.5), Inches(0.85),
             t, size=16, bold=True, color=AZUL_NAVY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.0)

# Placeholder render
add_rect(s, Inches(8.8), Inches(1.5), Inches(4.0), Inches(4.5), AZUL_NAVY)
add_text(s, Inches(8.8), Inches(3.4), Inches(4.0), Inches(0.5),
         '[RENDER TORRE]', size=14, color=BRANCO, align=PP_ALIGN.CENTER)

# Rodapé azul
add_rect(s, 0, Inches(6.5), W, Inches(1.0), AZUL_PRIM)
add_text(s, Inches(0.5), Inches(6.55), W - Inches(1.0), Inches(0.45),
         'ARBORIS — ORÇAMENTO PARAMÉTRICO V4', size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(7.0), W - Inches(1.0), Inches(0.35),
         'Arthen Empreendimentos  |  Itajaí/SC  |  Abril/2026',
         size=12, color=BRANCO, align=PP_ALIGN.CENTER)


# ============ SALVAR ============
import os
output_path = 'C:/Users/leona/orcamentos/parametricos/arthen-arboris/arthen-arboris-apresentacao-v00.pptx'
prs.save(output_path)
print(f"\n✅ Apresentação salva: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Tamanho: {os.path.getsize(output_path) / 1024:.0f} KB")
