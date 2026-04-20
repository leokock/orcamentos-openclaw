"""
Gera apresentacao .pptx do Arboris paramétrico v00 (primeira entrega ao cliente).
Baseada no memorial descritivo + media de mercado Cartesian.
Comparativos anonimos (medias, nao nomes de obras).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math

AZUL_PRIM = RGBColor(0x1E, 0x52, 0xF5)
AZUL_NAVY = RGBColor(0x0F, 0x1C, 0x4A)
LARANJA = RGBColor(0xFF, 0x57, 0x22)
CINZA_FUNDO = RGBColor(0xF4, 0xF6, 0xF8)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_LINHA = RGBColor(0xE8, 0xEE, 0xF4)
CINZA_TXT = RGBColor(0x45, 0x4E, 0x5E)
PRETO = RGBColor(0x1A, 0x1A, 0x1A)
VERDE = RGBColor(0x2E, 0x8B, 0x57)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, x, y, w, h, fill_color, line=False, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = fill_color
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=PRETO, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_footer(slide):
    add_text(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
             'Cartesian', size=11, bold=True, color=AZUL_PRIM, align=PP_ALIGN.RIGHT)


# ============ DADOS V00 ============
AC = 12472.98
UR = 98
CUB = 3028.45
TOTAL_V00 = 38811650  # calculado após ajustes gerenciamento

# Macrogrupos v00 (ordem da planilha)
MACROS = [
    ('GERENCIAMENTO TÉCNICO E ADMINISTRATIVO', 4039260, 10.4, 324),
    ('MOVIMENTAÇÃO DE TERRA', 187095, 0.5, 15),
    ('INFRAESTRUTURA', 2237313, 5.8, 179),
    ('SUPRAESTRUTURA', 9074861, 23.4, 728),
    ('ALVENARIA', 1336777, 3.4, 107),
    ('INSTALAÇÕES ELÉTRICAS, HIDRÁULICAS, GLP E PREVENTIVAS', 4918221, 12.7, 394),
    ('EQUIPAMENTOS E SISTEMAS ESPECIAIS', 1490000, 3.8, 119),
    ('CLIMATIZAÇÃO', 315930, 0.8, 25),
    ('IMPERMEABILIZAÇÃO', 669838, 1.7, 54),
    ('REVESTIMENTOS INTERNOS DE PAREDE', 1182223, 3.0, 95),
    ('REVESTIMENTOS E ACABAMENTOS EM TETO', 905729, 2.3, 73),
    ('PISOS E PAVIMENTAÇÕES', 2293599, 5.9, 184),
    ('PINTURA INTERNA', 1928859, 5.0, 155),
    ('ESQUADRIAS, VIDROS E FERRAGENS', 4053719, 10.4, 325),
    ('LOUÇAS E METAIS', 191100, 0.5, 15),
    ('FACHADA', 1894634, 4.9, 152),
    ('SERVIÇOS COMPLEMENTARES', 1518921, 3.9, 122),
    ('IMPREVISTOS E CONTINGÊNCIAS', 573634, 1.5, 46),
]

# ============================================================
# SLIDE 1 — Capa Cartesian
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0), Inches(2.8), W, Inches(1.8),
         'Cartesian', size=96, bold=True, color=AZUL_PRIM, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.5), W, Inches(0.5),
         'Coordenação Técnica e Orçamentária', size=20, color=CINZA_TXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — Capa do Projeto
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(1.2), Inches(6.5), Inches(0.9),
         'Orçamento\nParamétrico', size=44, bold=True, color=AZUL_PRIM)
add_rect(s, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.04), AZUL_PRIM)
add_text(s, Inches(0.8), Inches(3.5), Inches(6.5), Inches(0.9),
         'ARBORIS', size=50, bold=True, color=AZUL_NAVY)
add_text(s, Inches(0.8), Inches(4.5), Inches(6.5), Inches(0.4),
         'ARTHEN EMPREENDIMENTOS', size=18, bold=True, color=CINZA_TXT)
add_text(s, Inches(0.8), Inches(5.1), Inches(6.5), Inches(0.4),
         'Itapema/SC (Morretes)  |  Data-base: Março/2026  |  CUB/SC: R$ 3.028,45/m²',
         size=13, color=CINZA_TXT)
add_text(s, Inches(0.8), Inches(5.5), Inches(6.5), Inches(0.4),
         'Abril/2026',
         size=11, color=LARANJA, bold=True)
add_rect(s, Inches(7.5), Inches(0.8), Inches(5.3), Inches(5.9), AZUL_NAVY)
add_text(s, Inches(7.5), Inches(3.3), Inches(5.3), Inches(1.0),
         '[RENDER 3D DA TORRE]', size=16, color=BRANCO, align=PP_ALIGN.CENTER)
add_footer(s)


# ============================================================
# SLIDE 3 — Jornada
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.7),
         'JORNADA DA REUNIÃO', size=30, bold=True, color=AZUL_PRIM)

topics = [
    ('01', 'Dados e Premissas do Projeto'),
    ('02', 'Origem das Premissas'),
    ('03', 'Visão Geral do Orçamento'),
    ('04', 'Composição por Macrogrupo'),
    ('05', 'Especificações (conforme Memorial)'),
    ('06', 'Comparativo com Média de Mercado'),
    ('07', 'Justificativas — Itens Acima da Média'),
    ('08', 'Custo por Área Privativa e Tipologia'),
    ('09', 'Planejamento Macro e Curva Financeira'),
]
y = Inches(1.3)
for num, title in topics:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.58), BRANCO)
    add_rect(s, Inches(0.8), y, Inches(0.1), Inches(0.58), AZUL_PRIM)
    add_text(s, Inches(1.1), y, Inches(1.3), Inches(0.58),
             num, size=22, bold=True, color=AZUL_PRIM, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.4), y, Inches(9.5), Inches(0.58),
             title, size=15, bold=True, color=AZUL_NAVY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.65)
add_footer(s)


# ============================================================
# SLIDE 4 — Características do Empreendimento
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.7),
         'CARACTERÍSTICAS DO EMPREENDIMENTO', size=28, bold=True, color=AZUL_PRIM)

caracteristicas = [
    ('Área Construída', '12.472,98 m²'),
    ('Área do Terreno', '1.008,00 m²'),
    ('Unidades Residenciais', '90'),
    ('Unidades Comerciais', '8'),
    ('Total de Unidades', '98'),
    ('Pavimentos', '24 (Térreo + G1+G2+G3 + Dif + 14 Tipo + Rooftop)'),
    ('Apartamentos/Andar Tipo', '6 (2 suítes + lavabo + living com cozinha + sacada c/ churrasq.)'),
    ('Elevadores', '2 (1 social + 1 emergência)'),
    ('Vagas de Estacionamento', '99'),
    ('Sistema Estrutural', 'Concreto armado, estacas hélice contínua'),
    ('Padrão (memorial)', 'Médio — Interesse Social'),
    ('Prazo de Execução', '36 meses'),
]
y = Inches(1.5)
for k, v in caracteristicas:
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.4), AZUL_NAVY)
    add_text(s, Inches(1.0), y, Inches(5.0), Inches(0.4),
             k, size=12, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(6.2), y, Inches(6.3), Inches(0.4),
             v, size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.44)
add_footer(s)


# ============================================================
# SLIDE 5 — Origem das Premissas (NOVO)
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.7),
         'ORIGEM DAS PREMISSAS', size=28, bold=True, color=AZUL_PRIM)
add_text(s, Inches(0.8), Inches(1.05), Inches(12), Inches(0.4),
         'Cada premissa do paramétrico foi extraída do memorial do cliente ou estimada com base na média de mercado Cartesian.',
         size=12, color=CINZA_TXT)

# 2 colunas
# Coluna esquerda — Memorial
add_rect(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.5), BRANCO)
add_rect(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(0.5), AZUL_PRIM)
add_text(s, Inches(1.0), Inches(1.6), Inches(5.8), Inches(0.5),
         '📋 PREMISSAS DO MEMORIAL DO CLIENTE', size=13, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)

memorial_items = [
    'Localização: Ruas 418/420, Itapema/SC',
    'Área construída, terreno e nº de unidades',
    '24 pavimentos, 2 elevadores, 99 vagas',
    'Padrão declarado: Médio (Interesse Social)',
    'SEM subsolo (tudo acima do nível do solo)',
    'Estrutura: concreto armado + hélice contínua',
    'Laje em concreto armado mista (fck mín. 20 MPa)',
    'Alvenaria: bloco cerâmico (tijolo furado)',
    'Fachada: TINTA PREMIUM ACRÍLICA com texturas (2-4 cores)',
    'Pisos: porcelanato 60×60 + laminado em suítes',
    'Paredes BWC: azulejo 30×45',
    'Cozinha: pintura PVA+acrílica, azulejo só parede molhada',
    'Forros: gesso acartonado ou reboco',
    'Esquadrias: alumínio eletrostático + soleiras em granito',
    'Louças: SÓ bacias sanitárias (Deca/Incepa/Celite)',
    'Itens não fornecidos: cubas, bancadas, metais, chuveiros, box, luminárias, ar-condicionado',
]
y = Inches(2.2)
for item in memorial_items:
    add_text(s, Inches(1.0), y, Inches(5.7), Inches(0.28),
             '• ' + item, size=9, color=PRETO)
    y += Inches(0.3)

# Coluna direita — Cartesian
add_rect(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.5), BRANCO)
add_rect(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(0.5), LARANJA)
add_text(s, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.5),
         '📊 ESTIMATIVAS CARTESIAN (média de mercado)', size=13, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)

cartesian_items = [
    'Prazo de execução: 36 meses (mediana torre 24 pav.)',
    'Gerenciamento: PUs de equipe técnica, vigilância, EPI, consumos',
    'PUs materiais de concreto, aço e forma',
    'Índices físicos (m²/m² AC, kg/m² AC, m³/m² AC)',
    'MO empreitada por disciplina',
    'Sistemas especiais: CFTV, interfonia, automação, bombas',
    'Impermeabilização: PU da manta e MO',
    'Pisos: contrapiso + rodapé + MO assentamento',
    'Pintura: selador + PVA + acrílica (PUs)',
    'Esquadrias: PU de alumínio + serralheria',
    'Imprevistos: 1,5% (padrão Cartesian para paramétrico)',
    'Complementares: paisagismo, mobiliário, ligações definitivas',
]
y = Inches(2.2)
for item in cartesian_items:
    add_text(s, Inches(7.2), y, Inches(5.4), Inches(0.28),
             '• ' + item, size=9, color=PRETO)
    y += Inches(0.3)

add_text(s, Inches(0.8), Inches(7.0), Inches(12), Inches(0.3),
         'Base Cartesian: 126 obras Médio-Alto e Alto da Grande Florianópolis e Litoral Norte SC (2024-2026), indexadas ao CUB/SC Mar/2026',
         size=9, color=CINZA_TXT, align=PP_ALIGN.CENTER)
add_footer(s)


# ============================================================
# SLIDE 6 — Visão Geral do Orçamento
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
         'VISÃO GERAL DO ORÇAMENTO', size=32, bold=True, color=AZUL_PRIM)

cards = [
    ('R$ 38,8 Mi', 'CUSTO TOTAL', 'Valor global do empreendimento'),
    ('R$ 3.112/m²', 'CUSTO POR M²', 'Área construída: 12.473 m²'),
    ('R$ 396 mil', 'CUSTO POR UNIDADE', 'Total de 98 unidades'),
    ('1,03 CUB', 'FATOR CUB/SC', 'CUB Mar/2026: R$ 3.028,45'),
]
xs = [Inches(0.8), Inches(7.0)]
ys = [Inches(1.6), Inches(4.3)]
for i, (big, label, desc) in enumerate(cards):
    x = xs[i % 2]
    y = ys[i // 2]
    add_rect(s, x, y, Inches(5.5), Inches(2.5), AZUL_NAVY)
    add_text(s, x, y+Inches(0.25), Inches(5.5), Inches(0.85),
             big, size=36, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
    add_text(s, x, y+Inches(1.25), Inches(5.5), Inches(0.4),
             label, size=14, bold=True, color=AZUL_PRIM, align=PP_ALIGN.CENTER)
    add_text(s, x, y+Inches(1.75), Inches(5.5), Inches(0.5),
             desc, size=13, color=BRANCO, align=PP_ALIGN.CENTER)
add_footer(s)


# ============================================================
# SLIDE 7 — Detalhamento do Custo
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'DETALHAMENTO DO CUSTO', size=28, bold=True, color=AZUL_PRIM)

total_v = sum(x[1] for x in MACROS)
hdr_y = Inches(1.2)
add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.4), AZUL_PRIM)
add_text(s, Inches(0.6), hdr_y, Inches(6.6), Inches(0.4), 'ETAPA',
         size=11, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.2), hdr_y, Inches(2.2), Inches(0.4), 'VALOR ORÇADO',
         size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.6), hdr_y, Inches(1.2), Inches(0.4), '%',
         size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.0), hdr_y, Inches(1.7), Inches(0.4), 'VALOR / m²',
         size=11, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

y = Inches(1.6)
for i, (name, val, pct, rsm2) in enumerate(MACROS):
    bg = BRANCO if i % 2 == 0 else CINZA_LINHA
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.27), bg)
    add_text(s, Inches(0.6), y, Inches(6.6), Inches(0.27),
             name, size=9, color=PRETO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.2), y, Inches(2.2), Inches(0.27),
             f'R$ {val:,.2f}'.replace(',', '.'), size=9, color=PRETO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.6), y, Inches(1.2), Inches(0.27),
             f'{pct:.1f}%', size=9, color=PRETO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(11.0), y, Inches(1.7), Inches(0.27),
             f'R$ {rsm2:,.2f}'.replace(',', '.'), size=9, color=PRETO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.27)

add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.35), AZUL_PRIM)
add_text(s, Inches(0.6), y, Inches(6.6), Inches(0.35), 'TOTAL',
         size=12, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.2), y, Inches(2.2), Inches(0.35),
         f'R$ {total_v:,.2f}'.replace(',', '.'), size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.6), y, Inches(1.2), Inches(0.35), '100%',
         size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.0), y, Inches(1.7), Inches(0.35),
         f'R$ {total_v/AC:,.2f}'.replace(',', '.'), size=12, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
add_footer(s)


# ============================================================
# SLIDE 8 — Composição (top 10 barras + TOP 3 cards)
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.8),
         'COMPOSIÇÃO DO ORÇAMENTO', size=30, bold=True, color=AZUL_PRIM)

shortnames = {
    'GERENCIAMENTO TÉCNICO E ADMINISTRATIVO': 'Gerenciamento',
    'SUPRAESTRUTURA': 'Supraestrutura',
    'INFRAESTRUTURA': 'Infraestrutura',
    'INSTALAÇÕES ELÉTRICAS, HIDRÁULICAS, GLP E PREVENTIVAS': 'Instalações',
    'ESQUADRIAS, VIDROS E FERRAGENS': 'Esquadrias',
    'FACHADA': 'Fachada',
    'ALVENARIA': 'Alvenaria',
    'SERVIÇOS COMPLEMENTARES': 'Complementares',
    'PISOS E PAVIMENTAÇÕES': 'Pisos',
    'PINTURA INTERNA': 'Pintura',
    'REVESTIMENTOS INTERNOS DE PAREDE': 'Rev. Parede',
    'EQUIPAMENTOS E SISTEMAS ESPECIAIS': 'Sist. Especiais',
    'REVESTIMENTOS E ACABAMENTOS EM TETO': 'Rev. Teto',
    'IMPERMEABILIZAÇÃO': 'Impermeabilização',
    'LOUÇAS E METAIS': 'Louças e Metais',
    'MOVIMENTAÇÃO DE TERRA': 'Mov. Terra',
    'IMPREVISTOS E CONTINGÊNCIAS': 'Imprevistos',
    'CLIMATIZAÇÃO': 'Climatização',
}

top10 = sorted([(shortnames.get(m[0], m[0]), m[2]) for m in MACROS], key=lambda x: -x[1])[:10]
top10 = list(reversed(top10))
max_pct = max(v for _, v in top10)
chart_x = Inches(0.8)
chart_y = Inches(1.4)
chart_w = Inches(7.3)
chart_h = Inches(5.7)
bar_h = Emu(int(int(chart_h - Inches(0.3)) / len(top10)))
for i, (name, v) in enumerate(top10):
    y = chart_y + bar_h * i
    add_text(s, chart_x, y, Inches(1.8), bar_h,
             name, size=11, color=CINZA_TXT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    bar_w = Emu(int(int(chart_w - Inches(2.0)) * v / max_pct))
    add_rect(s, chart_x + Inches(1.9), y + Emu(int(int(bar_h) * 0.2)),
             bar_w, Emu(int(int(bar_h) * 0.6)), AZUL_PRIM)
    add_text(s, chart_x + Inches(1.9) + bar_w + Inches(0.05), y,
             Inches(1.0), bar_h,
             f'{v:.1f}%', size=11, bold=True, color=AZUL_NAVY, anchor=MSO_ANCHOR.MIDDLE)

top3_data = sorted(MACROS, key=lambda x: -x[2])[:3]
card_x = Inches(8.4)
card_y = Inches(1.4)
card_w = Inches(4.6)
card_h = Inches(1.75)
reasons = {
    'SUPRAESTRUTURA': 'Laje convencional — torre de 24 pavimentos',
    'INSTALAÇÕES ELÉTRICAS, HIDRÁULICAS, GLP E PREVENTIVAS': 'Tipologia 2 suítes + lavabo (3 BWCs/apto)',
    'ESQUADRIAS, VIDROS E FERRAGENS': 'Alumínio eletrostático + granito soleiras (memorial)',
    'GERENCIAMENTO TÉCNICO E ADMINISTRATIVO': '36 meses de obra — equipe enxuta',
}
add_text(s, card_x, card_y - Inches(0.1), card_w, Inches(0.4),
         'TOP 3 Macrogrupos', size=16, bold=True, color=AZUL_NAVY)
for i, (name, val, pct, rsm2) in enumerate(top3_data):
    y = card_y + Inches(0.35) + card_h * i
    add_rect(s, card_x, y, card_w, card_h - Inches(0.15), BRANCO)
    add_rect(s, card_x, y, Inches(0.08), card_h - Inches(0.15), LARANJA)
    add_text(s, card_x + Inches(0.25), y + Inches(0.1), Inches(1.8), Inches(0.6),
             f'{pct:.1f}%', size=24, bold=True, color=AZUL_PRIM)
    short = shortnames.get(name, name)
    add_text(s, card_x + Inches(2.0), y + Inches(0.15), Inches(2.5), Inches(0.4),
             short, size=13, bold=True, color=AZUL_NAVY)
    add_text(s, card_x + Inches(0.25), y + Inches(0.75), Inches(1.8), Inches(0.35),
             f'R$ {val/1e6:.1f} Mi', size=14, bold=True, color=LARANJA)
    add_text(s, card_x + Inches(2.0), y + Inches(0.75), Inches(2.5), Inches(0.75),
             reasons.get(name, '—'), size=10, color=CINZA_TXT)
add_footer(s)


# ============================================================
# SLIDES 9-12 — Especificações (conforme memorial)
# ============================================================
specs_blocos = [
    # Slide 9
    [
        ('SUPRAESTRUTURA', '22,8%', 'R$ 9.074.861  |  R$ 728/m²',
         ['Laje convencional em concreto armado (memorial: mista conforme projeto estrutural)',
          'fck 30 MPa + aço CA-50/CA-60 (memorial)',
          'Formas em compensado plastificado 18mm']),
        ('GERENCIAMENTO TÉCNICO', '10,4%', 'R$ 4.039.260  |  R$ 324/m²',
         ['Prazo de 36 meses de obra com equipe enxuta',
          'Equipe: engenheiro, mestre, encarregado, almoxarife, vigilância, limpeza',
          'Projetos + consultorias + taxas + ensaios + canteiro + cremalheira']),
        ('INSTALAÇÕES', '12,3%', 'R$ 4.918.221  |  R$ 394/m²',
         ['Elétrica, hidro, preventivas (hidrantes + alarme + SPDA), gás GLP e telecom',
          'Captação de água da chuva + cisterna (memorial)',
          'Pontos de ar-condicionado split em suítes e sala (memorial — só infra)']),
        ('ESQUADRIAS, VIDROS E FERRAGENS', '10,2%', 'R$ 4.053.719  |  R$ 325/m²',
         ['Janelas de alumínio com pintura eletrostática preta (memorial)',
          'Portas laminadas semi-ocas brancas (memorial)',
          'Ferragens Papaiz/Arouca/Pado + soleiras em granito']),
    ],
    # Slide 10
    [
        ('PISOS E PAVIMENTAÇÕES', '5,8%', 'R$ 2.293.599  |  R$ 184/m²',
         ['Sala, circulação, BWCs, cozinha, A.S. e lavabos: PORCELANATO 60×60 (memorial)',
          'Suítes: piso LAMINADO Eucafloor/Durafloor (memorial)',
          'Estacionamento: cimento alisado. Halls comuns: granito polido']),
        ('INFRAESTRUTURA', '5,6%', 'R$ 2.237.313  |  R$ 179/m²',
         ['Fundação em estacas hélice contínua (memorial — NBR 6122)',
          'Blocos, baldrames, contenção e infraestrutura de serviço',
          'Contrata MO de fundação em regime empreitada']),
        ('PINTURA INTERNA', '4,8%', 'R$ 1.928.859  |  R$ 155/m²',
         ['Áreas privativas: pintura PVA sobre acrílica (memorial — Suvinil/Coral)',
          'Estacionamento e áreas técnicas: textura ou látex',
          'Demarcação epóxi de vagas conforme norma']),
        ('FACHADA', '4,8%', 'R$ 1.894.634  |  R$ 152/m²',
         ['Tinta premium acrílica com texturas em 2 a 4 cores (memorial)',
          'Muros: tinta premium acrílica lisa ou textura',
          'Balancim fachadeiro (torre vertical de 24 pavimentos)']),
    ],
    # Slide 11
    [
        ('COMPLEMENTARES', '3,8%', 'R$ 1.518.921  |  R$ 122/m²',
         ['Mobiliário e decoração de áreas comuns + paisagismo',
          'Rooftop com lazer completo: 2 piscinas, ofurô, academia, gourmets, coworking',
          'Limpeza final + ligações definitivas + desmobilização']),
        ('EQUIPAMENTOS E SIST. ESPECIAIS', '3,7%', 'R$ 1.490.000  |  R$ 119/m²',
         ['2 elevadores (1 social + 1 emergência) — Atlas/Thyssen/Otis (memorial)',
          'Gerador dedicado + CFTV + interfonia + automação (BMS)',
          'SPDA (para-raios) + bombas de recalque + quadros de comando']),
        ('ALVENARIA', '3,4%', 'R$ 1.336.777  |  R$ 107/m²',
         ['Bloco cerâmico 14cm — tijolo furado em todas as vedações (memorial)',
          'Vergas e contravergas em concreto armado',
          'Bloco de concreto celular na escada enclausurada (resistência ao fogo)']),
        ('REV. INTERNO DE PAREDE', '3,0%', 'R$ 1.182.223  |  R$ 95/m²',
         ['Reboco massa única + chapisco em todas as áreas',
          'Azulejo 30×45 em banheiros e parede molhada de cozinha (memorial)',
          'Porcelanato em halls principais do edifício']),
    ],
    # Slide 12
    [
        ('REV. E ACABAMENTOS EM TETO', '2,3%', 'R$ 905.729  |  R$ 73/m²',
         ['Forro de gesso acartonado liso em áreas privativas (memorial)',
          'Reboco + pintura acrílica em áreas comuns de estacionamento',
          'Forro RU (resistente à umidade) em BWCs e áreas molhadas']),
        ('IMPERMEABILIZAÇÃO', '1,7%', 'R$ 669.838  |  R$ 54/m²',
         ['Manta asfáltica 4mm + argamassa polimérica (memorial — conforme normas vigentes)',
          'Vigas de fundação + BWCs + sacadas + terraços + caixa d\'água + lajes',
          'Piscinas + forro casa de máquinas (memorial)']),
        ('IMPREVISTOS E CONTINGÊNCIAS', '1,5%', 'R$ 589.156  |  R$ 47/m²',
         ['Percentual aplicado sobre o subtotal do orçamento',
          'Padrão Cartesian para orçamentos pré-executivos',
          'Cobre desvios menores na transição para executivo']),
        ('CLIMATIZAÇÃO + LOUÇAS + MOV. TERRA', '1,8%', 'R$ 694.125  |  R$ 55/m²',
         ['Climatização: só infraestrutura (dreno + eletroduto) — aparelhos não fornecidos',
          'Louças: só bacias sanitárias Deca/Incepa/Celite (itens não fornecidos pelo cliente)',
          'Mov. Terra: escavação de fundações (sem subsolo — memorial)']),
    ],
]

for block in specs_blocos:
    s = add_blank_slide()
    add_rect(s, 0, 0, W, H, CINZA_FUNDO)
    add_text(s, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
             'ESPECIFICAÇÕES (conforme memorial do cliente)', size=26, bold=True, color=AZUL_PRIM)
    y = Inches(1.3)
    for title, pct, val, bullets in block:
        add_rect(s, Inches(0.8), y, Inches(0.1), Inches(1.35), LARANJA)
        add_rect(s, Inches(0.9), y, Inches(11.7), Inches(1.35), BRANCO)
        add_text(s, Inches(1.0), y+Inches(0.05), Inches(6.5), Inches(0.4),
                 title, size=13, bold=True, color=AZUL_NAVY)
        add_text(s, Inches(7.3), y+Inches(0.05), Inches(1.5), Inches(0.4),
                 pct, size=15, bold=True, color=AZUL_PRIM, align=PP_ALIGN.RIGHT)
        add_text(s, Inches(8.9), y+Inches(0.05), Inches(3.6), Inches(0.4),
                 val, size=11, color=CINZA_TXT, align=PP_ALIGN.RIGHT)
        by = y + Inches(0.5)
        for b in bullets:
            add_text(s, Inches(1.0), by, Inches(11.5), Inches(0.3),
                     '• ' + b, size=10, color=PRETO)
            by += Inches(0.28)
        y += Inches(1.45)
    add_footer(s)


# ============================================================
# SLIDE 13 — Comparativo com Média de Mercado (M-A)
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'COMPARATIVO COM MÉDIA DE MERCADO', size=28, bold=True, color=AZUL_PRIM)
add_text(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
         'Referência: média de 60 obras padrão Médio-Alto da base Cartesian (Grande Florianópolis e Litoral Norte SC, 2024-2026)',
         size=11, color=CINZA_TXT)

comp = [
    ('Gerenciamento', 324, 475, '-32%'),
    ('Mov. Terra', 15, 115, '-87%'),
    ('Infraestrutura', 179, 196, '-9%'),
    ('Supraestrutura', 728, 655, '+11%'),
    ('Alvenaria', 107, 140, '-24%'),
    ('Instalações', 394, 332, '+19%'),
    ('Sist. Especiais', 119, 173, '-31%'),
    ('Climatização', 25, 31, '-19%'),
    ('Impermeabilização', 54, 60, '-10%'),
    ('Rev. Parede', 95, 61, '+56%*'),
    ('Teto', 73, 60, '+22%'),
    ('Pisos', 184, 199, '-8%'),
    ('Pintura', 155, 125, '+24%'),
    ('Esquadrias', 325, 301, '+8%'),
    ('Louças e Metais', 15, 60, '-75%'),
    ('Fachada', 152, 117, '+29%'),
    ('Complementares', 122, 201, '-39%'),
    ('Imprevistos', 46, 50, '-8%'),
    ('TOTAL', 3112, 3349, '-7,1%'),
]
hdr_y = Inches(1.6)
add_rect(s, Inches(1.8), hdr_y, Inches(9.7), Inches(0.4), AZUL_PRIM)
heads = [('ETAPA', 1.9, 3.5, PP_ALIGN.LEFT),
         ('ARBORIS', 5.4, 2.0, PP_ALIGN.CENTER),
         ('MÉDIA M-A', 7.4, 2.0, PP_ALIGN.CENTER),
         ('VARIAÇÃO', 9.4, 2.0, PP_ALIGN.CENTER)]
for txt, x, w, al in heads:
    add_text(s, Inches(x), hdr_y, Inches(w), Inches(0.4),
             txt, size=11, bold=True, color=BRANCO, align=al, anchor=MSO_ANCHOR.MIDDLE)

y = Inches(2.0)
for i, (name, arb, med, var) in enumerate(comp):
    is_total = name == 'TOTAL'
    bg = AZUL_PRIM if is_total else (BRANCO if i % 2 == 0 else CINZA_LINHA)
    txt_color = BRANCO if is_total else PRETO
    add_rect(s, Inches(1.8), y, Inches(9.7), Inches(0.23), bg)
    add_text(s, Inches(1.9), y, Inches(3.5), Inches(0.23),
             name, size=10, color=txt_color, bold=is_total, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.4), y, Inches(2.0), Inches(0.23),
             f'R$ {arb:,}'.replace(',', '.'), size=10, color=txt_color, bold=is_total, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.4), y, Inches(2.0), Inches(0.23),
             f'R$ {med:,}'.replace(',', '.'), size=10, color=txt_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    var_color = BRANCO if is_total else (AZUL_PRIM if var.startswith('-') else LARANJA)
    add_text(s, Inches(9.4), y, Inches(2.0), Inches(0.23),
             var, size=10, bold=True, color=var_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.23)

add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
         '* Rev. Parede: amostra pequena na base (n=15), valores inconsistentes entre obras. Média não é estável.',
         size=9, color=CINZA_TXT)
add_text(s, Inches(0.5), Inches(7.15), Inches(12), Inches(0.3),
         'Valores em R$/m² AC, indexados ao CUB/SC Mar/2026.',
         size=9, color=CINZA_TXT)
add_footer(s)


# ============================================================
# SLIDES 14-15 — Justificativas itens acima da média (2 slides)
# ============================================================
justificativas = [
    ('SUPRAESTRUTURA', 728, 655, '+11%',
     'Torre de 24 pavimentos com laje maciça convencional',
     [
         'Altura: 24 pavimentos (Térreo + G1 + G2 + G3 + Diferenciado + 14 Tipo + Rooftop + Cobertura)',
         'Laje convencional maciça amplia consumo de concreto em ~25% vs nervurada',
         'Pé-direito padrão 3,00m em 19 pav. repetição amplifica dimensionamento estrutural',
         'Índices físicos na mediana da base (concreto 0,25 m³/m² AC, aço 106 kg/m² AC)',
     ]),
    ('INSTALAÇÕES', 394, 332, '+19%',
     'Tipologia 2 suítes + lavabo + cozinha (3 BWCs/apto, memorial)',
     [
         'Tipologia 2 suítes + lavabo eleva fator hidráulico em +15 a +20%',
         '12 pontos de água fria/quente por apto (cozinha, BWCs, área de serviço)',
         'Sistema preventivo completo: hidrantes, alarme, SPDA, para-raios (memorial)',
         'Captação de água da chuva + cisterna + medidor individual por apto',
     ]),
    ('ESQUADRIAS, VIDROS E FERRAGENS', 325, 301, '+8%',
     'Alumínio com pintura eletrostática + soleiras em granito (memorial)',
     [
         'Janelas de alumínio com pintura eletrostática preta (padrão superior a anodizado)',
         'Soleiras em granito em todas as janelas (memorial 5.7)',
         'Guarda-corpo de sacadas: alvenaria + esquadria alumínio envidraçada',
         'Ferragens zamac Papaiz/Arouca/Pado (memorial 5.9)',
     ]),
    ('PINTURA INTERNA', 155, 125, '+24%',
     'Entrega completa: PVA + acrílica em todas as áreas (memorial)',
     [
         'Áreas privativas: pintura PVA sobre acrílica (memorial — Suvinil/Coral)',
         'Estacionamento + escadas + áreas técnicas: textura ou látex completa',
         'Demarcação epóxi de vagas e pintura de pisos técnicos',
         'MO empreitada representa 66,2% do custo — mão de obra especializada',
     ]),
    ('FACHADA', 152, 117, '+29%',
     'Textura premium acrílica em torre vertical de 24 pavimentos (memorial)',
     [
         'Tinta premium acrílica com 2 a 4 cores e texturas (memorial 9.1)',
         'Muros: tinta premium acrílica lisa ou textura',
         'Índice fachada/AC de 1,55 m²/m² — normal para torre de 24 pav. (obras baixas: 0,8-1,2)',
         'PU conservador: R$ 98/m² de fachada (material + MO + balancim)',
     ]),
]

# Divide em 2 slides: 3 + 2
grupos = [justificativas[:3], justificativas[3:]]
for idx, grupo in enumerate(grupos):
    s = add_blank_slide()
    add_rect(s, 0, 0, W, H, CINZA_FUNDO)
    add_text(s, Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.6),
             f'JUSTIFICATIVAS — ITENS ACIMA DA MÉDIA ({idx+1}/{len(grupos)})',
             size=24, bold=True, color=AZUL_PRIM)
    add_text(s, Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.35),
             'Macrogrupos com valor superior à média de mercado e as razões técnicas',
             size=11, color=CINZA_TXT)

    y = Inches(1.5)
    block_h = Inches(1.85)
    for mg, arb, media, delta, subtitle, bullets in grupo:
        add_rect(s, Inches(0.5), y, Inches(12.3), block_h, BRANCO)
        add_rect(s, Inches(0.5), y, Inches(0.12), block_h, LARANJA)
        add_text(s, Inches(0.75), y + Inches(0.05), Inches(6.5), Inches(0.4),
                 mg, size=13, bold=True, color=AZUL_NAVY)
        add_rect(s, Inches(7.2), y + Inches(0.08), Inches(5.3), Inches(0.38), CINZA_FUNDO)
        add_text(s, Inches(7.3), y + Inches(0.08), Inches(2.5), Inches(0.38),
                 f'Arboris: R$ {arb}/m²', size=10, bold=True, color=AZUL_PRIM, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.8), y + Inches(0.08), Inches(1.6), Inches(0.38),
                 f'Média: R$ {media}/m²', size=10, color=CINZA_TXT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(11.4), y + Inches(0.08), Inches(1.05), Inches(0.38),
                 delta, size=11, bold=True, color=LARANJA, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(0.75), y + Inches(0.5), Inches(12), Inches(0.3),
                 subtitle, size=11, bold=True, color=AZUL_PRIM)
        by = y + Inches(0.85)
        for b in bullets[:3]:
            add_text(s, Inches(0.85), by, Inches(11.8), Inches(0.32),
                     '• ' + b, size=9, color=PRETO)
            by += Inches(0.32)
        y += block_h + Inches(0.1)
    add_footer(s)


# ============================================================
# SLIDE 16 — Custo por Área Privativa
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'ANÁLISE DE CUSTO POR ÁREA PRIVATIVA', size=28, bold=True, color=AZUL_PRIM)

# AP estimada: ~57% da AC (torres residenciais típicas)
AP_ARBORIS = round(AC * 0.57, 0)
rsm2_priv = TOTAL_V00 / AP_ARBORIS

priv_data = [
    (f'ARBORIS (este)', f'~{AP_ARBORIS:,.0f} m² (estimado)'.replace(',', '.'),
     f'R$ {rsm2_priv:,.0f}'.replace(',', '.'), True),
    ('Média Médio-Alto (n=37)', '—', 'R$ 5.650', False),
    ('Média Alto (n=23)', '—', 'R$ 7.100', False),
]

tab_y = Inches(1.7)
add_rect(s, Inches(1.5), tab_y, Inches(10.3), Inches(0.45), AZUL_NAVY)
add_text(s, Inches(1.7), tab_y, Inches(3.8), Inches(0.45),
         'EMPREENDIMENTO', size=13, bold=True, color=BRANCO, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.5), tab_y, Inches(3.5), Inches(0.45),
         'ÁREA PRIVATIVA', size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(9.0), tab_y, Inches(2.8), Inches(0.45),
         'R$/m² PRIV.', size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

y = tab_y + Inches(0.5)
for name, ap, rsm2, highlight in priv_data:
    bg = CINZA_LINHA if highlight else BRANCO
    txt_color = AZUL_PRIM if highlight else PRETO
    add_rect(s, Inches(1.5), y, Inches(10.3), Inches(0.5), bg)
    add_text(s, Inches(1.7), y, Inches(3.8), Inches(0.5),
             name, size=13, bold=highlight, color=txt_color, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.5), y, Inches(3.5), Inches(0.5),
             ap, size=12, color=PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.0), y, Inches(2.8), Inches(0.5),
             rsm2, size=14, bold=True, color=AZUL_PRIM if highlight else PRETO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)

# Gráfico
chart_y = Inches(4.5)
vals = [int(rsm2_priv), 5650, 7100]
labels = ['ARBORIS', 'Média M-A', 'Média Alto']
maxv = max(vals)
chart_x = Inches(2.7)
bar_w = Inches(2.2)
for i, (v, lbl) in enumerate(zip(vals, labels)):
    x = chart_x + Inches(i*2.6)
    bh = Inches(2.0 * v / maxv)
    y_bar = chart_y + Inches(2.0) - bh
    add_rect(s, x, y_bar, bar_w, bh, AZUL_PRIM if i == 0 else CINZA_TXT)
    add_text(s, x, y_bar - Inches(0.3), bar_w, Inches(0.3),
             f'R$ {v:,}'.replace(',', '.'), size=12, bold=True, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chart_y + Inches(2.1), bar_w, Inches(0.3),
             lbl, size=11, bold=(i==0), color=PRETO, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         'Valores de referência. Área privativa final depende do levantamento arquitetônico executivo.',
         size=9, color=CINZA_TXT)
add_footer(s)


# ============================================================
# SLIDE 17 — Planejamento Macro (Gantt 36 meses)
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'PLANEJAMENTO MACRO', size=28, bold=True, color=AZUL_PRIM)

N_MONTHS = 36
gantt = [
    ('Mov. Terra',          [1, 3]),
    ('Infraestrutura',      [2, 7]),
    ('Supraestrutura',      [5, 17]),
    ('Alvenaria',           [8, 22]),
    ('Instalações',         [9, 25]),
    ('Sist. Especiais',     [14, 26]),
    ('Climatização',        [14, 28]),
    ('Impermeabilização',   [12, 27]),
    ('Rev. Piso+Parede',    [14, 30]),
    ('Rev. Teto',           [15, 28]),
    ('Pintura',             [16, 32]),
    ('Esquadrias',          [18, 30]),
    ('Fachada',             [14, 28]),
    ('Complementares',      [28, 36]),
]

table_x = Inches(0.5)
table_y = Inches(1.3)
table_w = Inches(12.3)
label_w = Inches(2.5)
months_w = table_w - label_w
month_w = Emu(int(int(months_w) / N_MONTHS))
row_h = Inches(0.3)

for m in range(N_MONTHS):
    x = table_x + label_w + month_w * m
    add_rect(s, x, table_y, month_w, row_h, AZUL_PRIM)
    add_text(s, x, table_y, month_w, row_h,
             f'M{m+1}', size=7, bold=True, color=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

for i, (name, (start, end)) in enumerate(gantt):
    y = table_y + row_h + row_h * i
    add_rect(s, table_x, y, label_w, row_h, CINZA_LINHA if i % 2 == 0 else BRANCO)
    add_text(s, table_x + Inches(0.05), y, label_w, row_h,
             name, size=8, color=PRETO, anchor=MSO_ANCHOR.MIDDLE)
    for m in range(N_MONTHS):
        x = table_x + label_w + month_w * m
        if start <= m+1 <= end:
            add_rect(s, x, y, month_w, row_h, AZUL_PRIM)
        else:
            add_rect(s, x, y, month_w, row_h, BRANCO)

add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         'Cronograma macro indicativo — 36 meses de execução',
         size=9, color=CINZA_TXT)
add_footer(s)


# ============================================================
# SLIDE 18 — Curva de Distribuição Mensal
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'CURVA DE DISTRIBUIÇÃO DE CUSTOS', size=28, bold=True, color=AZUL_PRIM)

pesos = []
for m in range(1, N_MONTHS+1):
    p = math.exp(-((m - 17)**2) / (2 * 7**2))
    pesos.append(p)
soma = sum(pesos)
pct_dist = [p / soma * 100 for p in pesos]
valores = [p / 100 * TOTAL_V00 for p in pct_dist]

chart1_y = Inches(1.4)
chart1_x = Inches(0.5)
chart1_w = Inches(12.3)
bar_w2 = Emu(int(int(chart1_w - Inches(0.5)) / N_MONTHS))
maxp = max(pct_dist)
for m in range(N_MONTHS):
    v = pct_dist[m]
    bh = Inches(2.3 * v / maxp)
    x = chart1_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart1_y + Inches(2.4) - bh, bar_w2 - Emu(50000), bh, AZUL_PRIM)
    add_text(s, x - Emu(30000), chart1_y + Inches(2.4) - bh - Inches(0.25), bar_w2 + Emu(60000), Inches(0.2),
             f'{v:.1f}%', size=6, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chart1_y + Inches(2.45), bar_w2, Inches(0.15),
             f'M{m+1}', size=6, color=CINZA_TXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3),
         'Distribuição mensal (%)', size=12, bold=True, color=AZUL_NAVY)

chart2_y = Inches(4.4)
maxv = max(valores)
for m in range(N_MONTHS):
    v = valores[m]
    bh = Inches(2.0 * v / maxv)
    x = chart1_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart2_y + Inches(2.0) - bh, bar_w2 - Emu(50000), bh, LARANJA)
    add_text(s, x - Emu(30000), chart2_y + Inches(2.0) - bh - Inches(0.22), bar_w2 + Emu(60000), Inches(0.2),
             f'{v/1e6:.1f}', size=6, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.1), Inches(5), Inches(0.3),
         'Valor mensal (R$ Mi)', size=12, bold=True, color=AZUL_NAVY)
add_footer(s)


# ============================================================
# SLIDE 19 — Curva Acumulada
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
         'CURVA ACUMULADA DE CUSTOS', size=28, bold=True, color=AZUL_PRIM)

pct_acum = []
tot = 0
for p in pct_dist:
    tot += p
    pct_acum.append(tot)
val_acum = [p / 100 * TOTAL_V00 for p in pct_acum]

chart_y = Inches(1.4)
chart_x = Inches(0.5)
chart_w = Inches(12.3)
bar_w2 = Emu(int(int(chart_w - Inches(0.5)) / N_MONTHS))
for m in range(N_MONTHS):
    v = pct_acum[m]
    bh = Inches(2.3 * v / 100)
    x = chart_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart_y + Inches(2.4) - bh, bar_w2 - Emu(30000), bh, AZUL_PRIM)
    if m % 3 == 0:
        add_text(s, x - Emu(30000), chart_y + Inches(2.4) - bh - Inches(0.25),
                 bar_w2 + Emu(60000), Inches(0.2),
                 f'{v:.0f}%', size=7, bold=True, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, x, chart_y + Inches(2.45), bar_w2, Inches(0.15),
             f'M{m+1}', size=6, color=CINZA_TXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(1.1), Inches(5), Inches(0.3),
         'Acumulado (%)', size=12, bold=True, color=AZUL_NAVY)

chart2_y = Inches(4.4)
for m in range(N_MONTHS):
    v = val_acum[m]
    bh = Inches(2.0 * v / TOTAL_V00)
    x = chart_x + Inches(0.4) + bar_w2 * m
    add_rect(s, x, chart2_y + Inches(2.0) - bh, bar_w2 - Emu(30000), bh, LARANJA)
    if m % 6 == 0 or m == N_MONTHS - 1:
        add_text(s, x - Emu(40000), chart2_y + Inches(2.0) - bh - Inches(0.25),
                 bar_w2 + Emu(80000), Inches(0.2),
                 f'{v/1e6:.0f}', size=7, bold=True, color=AZUL_NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.1), Inches(5), Inches(0.3),
         'Acumulado (R$ Mi)', size=12, bold=True, color=AZUL_NAVY)
add_footer(s)


# ============================================================
# SLIDE 20 — Entregáveis
# ============================================================
s = add_blank_slide()
add_rect(s, 0, 0, W, H, CINZA_FUNDO)
add_text(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.8),
         'ENTREGÁVEIS', size=32, bold=True, color=AZUL_PRIM)

entregaveis = [
    ('01', 'Apresentação'),
    ('02', 'Orçamento Analítico (Excel)'),
    ('03', 'Planejamento Macro e Distribuição Físico-Financeira'),
]
y = Inches(2.0)
for num, t in entregaveis:
    add_rect(s, Inches(0.8), y, Inches(7.5), Inches(0.85), BRANCO)
    add_rect(s, Inches(0.8), y, Inches(0.12), Inches(0.85), AZUL_PRIM)
    add_text(s, Inches(1.1), y, Inches(1.5), Inches(0.85),
             num, size=30, bold=True, color=AZUL_PRIM, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), y, Inches(5.5), Inches(0.85),
             t, size=16, bold=True, color=AZUL_NAVY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.0)

add_rect(s, Inches(8.8), Inches(1.5), Inches(4.0), Inches(4.5), AZUL_NAVY)
add_text(s, Inches(8.8), Inches(3.4), Inches(4.0), Inches(0.5),
         '[RENDER TORRE]', size=14, color=BRANCO, align=PP_ALIGN.CENTER)

add_rect(s, 0, Inches(6.5), W, Inches(1.0), AZUL_PRIM)
add_text(s, Inches(0.5), Inches(6.55), W - Inches(1.0), Inches(0.45),
         'ARBORIS — ORÇAMENTO PARAMÉTRICO', size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(7.0), W - Inches(1.0), Inches(0.35),
         'Arthen Empreendimentos  |  Itapema/SC  |  Abril/2026',
         size=12, color=BRANCO, align=PP_ALIGN.CENTER)


# ============ SALVAR ============
import os
output_path = 'C:/Users/leona/orcamentos/parametricos/arthen-arboris/arthen-arboris-apresentacao-v00b.pptx'
prs.save(output_path)
print(f"\n✅ Apresentacao v00 salva: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Tamanho: {os.path.getsize(output_path) / 1024:.0f} KB")
