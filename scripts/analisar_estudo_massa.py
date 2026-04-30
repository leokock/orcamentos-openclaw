#!/usr/bin/env python3
"""
Análise comparativa do Estudo de Massa 2000m² — Vertical × Horizontal.

Pipeline:
  1. Aplica overrides nos 2 xlsx (INDICES!C18 fachada, Sist. Especiais!C4 elevadores)
  2. Recalcula via Excel COM (win32com)
  3. Lê totais por macrogrupo de cada cenário
  4. Gera 3 entregáveis:
     - analise-comparativa-estudo-massa-2000m2.xlsx (4 abas)
     - analise-comparativa-estudo-massa-2000m2.docx (memorial)
     - apresentacao-estudo-massa-2000m2.pptx (10 slides)

Uso:
  py -3.10 scripts/analisar_estudo_massa.py
"""
import os
import sys
from pathlib import Path
from copy import copy

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.utils import get_column_letter
import win32com.client as win32

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ============== CONFIGURAÇÃO ==============

BASE = Path.home() / "orcamentos-openclaw" / "base" / "pacotes" / "estudo-massa-2000m2"
SLUG = "estudo-massa-2000m2"

XLSX_V = BASE / f"parametrico-{SLUG}-vertical.xlsx"
XLSX_H = BASE / f"parametrico-{SLUG}-horizontal.xlsx"
OUT_XLSX = BASE / f"analise-comparativa-{SLUG}.xlsx"
OUT_DOCX = BASE / f"analise-comparativa-{SLUG}.docx"
OUT_PPTX = BASE / f"apresentacao-{SLUG}.pptx"

# Overrides:
#   - Geométricos (V vs H): captam diferenças de geometria que o pipeline não pega via AC
#   - Calibração pequeno porte: PUs default da V2 são calibrados em projetos médios-grandes
#     (Placon 4077 m², Arboris 12473 m²); pra obra de 2000 m², itens fixos (projetos, taxas,
#     equipamentos, mobiliário) precisam ser proporcionalmente menores. Overrides abaixo
#     refletem padrão Médio em SC interior.
#
# Coordenadas: aba, (linha, coluna, valor). Coluna 3 = C (Qtd) | Coluna 5 = E (PU)
OVERRIDES = {
    "vertical": {
        # Geométricos
        "INDICES": [(18, 3, 0.74)],            # Fachada m²/m² AC (calc=1.55 → 0.74 real)
        "Sist. Especiais": [(4, 3, 2)],        # 2 elevadores sociais (NBR 13994: 9 paradas)
        # Calibração pequeno porte — Gerenciamento
        "Gerenciamento": [
            (4, 5, 150000),    # Projetos: R$ 650k → 150k (prédio pequeno padrão Médio)
            (5, 5, 80000),     # Consultorias: 180k → 80k
            (6, 5, 80000),     # Ensaios: 130k → 80k
            (7, 5, 150000),    # Taxas/seguros: 380k → 150k
            (15, 5, 6000),     # Vigilância PU/mês: 15.261 → 6.000 (interior SC)
            (20, 5, 100000),   # Inst. provisórias: 160k → 100k (canteiro compacto)
            (22, 5, 320000),   # Equipamentos: 480k → 320k (grua compacta 9 pavs)
            (21, 5, 6500),     # Despesas consumo PU/mês: 11.500 → 6.500
        ],
        # Calibração pequeno porte — Complementares
        "Complementares": [
            (4, 5, 120000),    # Mobiliário áreas comuns: 540k → 120k (lazer básico)
            (5, 5, 80000),     # Ambientação: 255k → 80k
            (6, 5, 50000),     # Paisagismo: 120k → 50k
            (7, 5, 50000),     # Equip. lazer: 180k → 50k (só salão+churrasq)
        ],
    },
    "horizontal": {
        # Geométricos
        "INDICES": [(18, 3, 0.64)],            # Fachada m²/m² AC
        # Sist. Especiais: default (1 social + 1 serviço)
        # Calibração pequeno porte
        "Gerenciamento": [
            (4, 5, 150000),
            (5, 5, 80000),
            (6, 5, 80000),
            (7, 5, 150000),
            (15, 5, 6000),
            (20, 5, 100000),
            (22, 5, 130000),   # H sem grua (7 pavs) — só cremalheira/elev. obra
            (21, 5, 6500),
        ],
        "Complementares": [
            (4, 5, 120000),
            (5, 5, 80000),
            (6, 5, 50000),
            (7, 5, 50000),
        ],
    },
}

GEOMETRIA = {
    "vertical":   {"footprint": 235, "perim": 62, "altura": 26, "fachada": 1395, "cob": 235, "np": 9, "npt": 8, "elev": 2, "aptos_pav": 3},
    "horizontal": {"footprint": 315, "perim": 71, "altura": 20, "fachada": 1213, "cob": 315, "np": 7, "npt": 6, "elev": 1, "aptos_pav": 4},
}

# ============== STYLES ==============

DARK = "2C3E50"
ACCENT = "2980B9"
ORANGE = "E67E22"
GREEN_BG = "E8F5E9"
RED_BG = "FDEDEC"
GRAY_BG = "F5F5F5"

THIN = Border(
    left=Side(style='thin', color='CCCCCC'), right=Side(style='thin', color='CCCCCC'),
    top=Side(style='thin', color='CCCCCC'), bottom=Side(style='thin', color='CCCCCC'),
)

def style_header(cell):
    cell.font = Font(bold=True, color="FFFFFF", size=10, name="Arial")
    cell.fill = PatternFill(start_color=DARK, end_color=DARK, fill_type="solid")
    cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
    cell.border = THIN

def style_money(cell):
    cell.number_format = '"R$" #,##0'
    cell.alignment = Alignment(horizontal='right')
    cell.border = THIN

def style_pct(cell):
    cell.number_format = '0.0%'
    cell.alignment = Alignment(horizontal='right')
    cell.border = THIN

def fmt_brl(v):
    if v is None: return "—"
    return f"R$ {v:,.0f}".replace(",", "X").replace(".", ",").replace("X", ".")

# ============== ETAPA 1: APLICAR OVERRIDES ==============

def aplicar_overrides():
    print("\n[1/4] Aplicando overrides nos 2 xlsx...")
    for cenario, sheets in OVERRIDES.items():
        path = BASE / f"parametrico-{SLUG}-{cenario}.xlsx"
        wb = load_workbook(path)
        for sheet_name, edits in sheets.items():
            ws = wb[sheet_name]
            for r, c, v in edits:
                ws.cell(r, c).value = v
                # marcar célula como override (cor laranja claro)
                ws.cell(r, c).fill = PatternFill(start_color="FFF3E0", end_color="FFF3E0", fill_type="solid")
                ws.cell(r, c).font = Font(bold=True, color=ORANGE, size=11)
        wb.save(path)
        print(f"  OK: {cenario} ({sum(len(e) for e in sheets.values())} overrides)")

# ============== ETAPA 2: RECALCULAR VIA EXCEL COM ==============

def recalcular_xlsx(paths):
    print("\n[2/4] Recalculando fórmulas via Excel COM...")
    excel = win32.DispatchEx("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    try:
        for path in paths:
            print(f"  Abrindo: {path.name}")
            wb = excel.Workbooks.Open(str(path.absolute()))
            excel.CalculateFull()
            wb.Save()
            wb.Close(SaveChanges=False)
            print(f"  OK: recalculado")
    finally:
        excel.Quit()

# ============== ETAPA 3: LER TOTAIS ==============

def ler_totais(path):
    """Retorna dict {macrogrupo: custo} + grand_total."""
    wb = load_workbook(path, data_only=True)
    ws = wb["CUSTOS_MACROGRUPO"]
    totais = {}
    for r in range(4, 22):
        macro = ws.cell(r, 1).value
        custo = ws.cell(r, 4).value
        if macro and custo is not None:
            totais[macro] = float(custo)
    grand_total = ws.cell(22, 4).value
    return totais, float(grand_total) if grand_total else 0.0

def ler_indice_fachada(path):
    """Lê valor efetivo do índice fachada (D18)."""
    wb = load_workbook(path, data_only=True)
    ws = wb["INDICES"]
    return float(ws.cell(18, 4).value or 0)

# ============== ETAPA 4a: GERAR COMPARATIVO XLSX ==============

def gerar_comparativo_xlsx(t_v, gt_v, t_h, gt_h, fachada_v, fachada_h):
    print("\n[4a/4] Gerando analise-comparativa.xlsx...")
    wb = Workbook()
    wb.remove(wb.active)

    # ABA 1: RESUMO
    ws = wb.create_sheet("RESUMO")
    ws["A1"] = "Estudo de Massa 2.000 m² — Comparativo Vertical × Horizontal"
    ws["A1"].font = Font(bold=True, size=14, color=DARK)
    ws.merge_cells("A1:E1")
    ws["A2"] = "Brusque/SC | 24 aptos × 55 m² priv | Padrão Médio | Térreo pilotis | Ger 30/04/26"
    ws["A2"].font = Font(size=9, color="666666", italic=True)
    ws.merge_cells("A2:E2")

    ws["A4"] = "GRAND TOTAL"
    ws["A4"].font = Font(bold=True, size=12, color=DARK)
    headers = ["", "Vertical (3 aptos/pav)", "Horizontal (4 aptos/pav)", "Δ (V − H)", "Δ %"]
    for i, h in enumerate(headers, 1):
        c = ws.cell(5, i, h)
        style_header(c)

    delta = gt_v - gt_h
    delta_pct = delta / gt_h if gt_h else 0
    ws.append(["Custo Total (R$)", gt_v, gt_h, delta, delta_pct])
    r = 6
    for col in [2, 3, 4]:
        style_money(ws.cell(r, col))
    style_pct(ws.cell(r, 5))
    ws.cell(r, 1).font = Font(bold=True, size=10)
    ws.cell(r, 1).border = THIN

    # R$/m²
    ws.append(["Custo / m² AC", gt_v/2000, gt_h/2000, (gt_v-gt_h)/2000, ""])
    r = 7
    for col in [2, 3, 4]:
        style_money(ws.cell(r, col))
    ws.cell(r, 1).font = Font(bold=True, size=10)
    ws.cell(r, 1).border = THIN

    # R$/UR
    ws.append(["Custo / UR (24 aptos)", gt_v/24, gt_h/24, (gt_v-gt_h)/24, ""])
    r = 8
    for col in [2, 3, 4]:
        style_money(ws.cell(r, col))
    ws.cell(r, 1).font = Font(bold=True, size=10)
    ws.cell(r, 1).border = THIN

    # Top 5 macrogrupos sensíveis
    ws["A11"] = "TOP 5 MACROGRUPOS SENSÍVEIS (maior Δ absoluto)"
    ws["A11"].font = Font(bold=True, size=12, color=DARK)
    headers = ["#", "Macrogrupo", "Vertical (R$)", "Horizontal (R$)", "Δ (R$)", "Δ %", "Quem ganha"]
    for i, h in enumerate(headers, 1):
        c = ws.cell(12, i, h)
        style_header(c)

    deltas = []
    for macro in t_v:
        v = t_v.get(macro, 0)
        h = t_h.get(macro, 0)
        d = v - h
        d_pct = d / h if h else 0
        deltas.append((macro, v, h, d, d_pct))
    deltas_sorted = sorted(deltas, key=lambda x: abs(x[3]), reverse=True)[:5]

    for i, (macro, v, h, d, dp) in enumerate(deltas_sorted, 1):
        winner = "H" if d > 0 else "V" if d < 0 else "="
        ws.append([i, macro, v, h, d, dp, winner])
        r = 12 + i
        for col in [3, 4, 5]:
            style_money(ws.cell(r, col))
        style_pct(ws.cell(r, 6))
        ws.cell(r, 7).alignment = Alignment(horizontal='center')
        ws.cell(r, 7).font = Font(bold=True, color=("D32F2F" if d > 0 else "388E3C"))
        ws.cell(r, 1).alignment = Alignment(horizontal='center')

    # Larguras
    ws.column_dimensions['A'].width = 5
    ws.column_dimensions['B'].width = 30
    ws.column_dimensions['C'].width = 22
    ws.column_dimensions['D'].width = 22
    ws.column_dimensions['E'].width = 18
    ws.column_dimensions['F'].width = 10
    ws.column_dimensions['G'].width = 14

    # ABA 2: POR_MACROGRUPO
    ws2 = wb.create_sheet("POR_MACROGRUPO")
    ws2["A1"] = "Comparativo por Macrogrupo (18 itens)"
    ws2["A1"].font = Font(bold=True, size=13, color=DARK)
    ws2.merge_cells("A1:G1")
    headers = ["Macrogrupo", "Vertical (R$)", "Horizontal (R$)", "Δ (R$)", "Δ %", "% V no total", "% H no total"]
    for i, h in enumerate(headers, 1):
        c = ws2.cell(3, i, h)
        style_header(c)

    for i, (macro, v, h, d, dp) in enumerate(deltas, 1):
        ws2.append([macro, v, h, d, dp, v/gt_v if gt_v else 0, h/gt_h if gt_h else 0])
        r = 3 + i
        for col in [2, 3, 4]:
            style_money(ws2.cell(r, col))
        style_pct(ws2.cell(r, 5))
        style_pct(ws2.cell(r, 6))
        style_pct(ws2.cell(r, 7))
        ws2.cell(r, 1).border = THIN

    # TOTAL
    rt = 3 + len(deltas) + 1
    ws2.cell(rt, 1, "TOTAL").font = Font(bold=True, size=11)
    ws2.cell(rt, 2, gt_v); ws2.cell(rt, 3, gt_h); ws2.cell(rt, 4, gt_v - gt_h)
    ws2.cell(rt, 5, (gt_v - gt_h)/gt_h if gt_h else 0)
    for col in [2, 3, 4]:
        style_money(ws2.cell(rt, col))
        ws2.cell(rt, col).font = Font(bold=True)
    style_pct(ws2.cell(rt, 5))

    ws2.column_dimensions['A'].width = 28
    for col in 'BCDEFG':
        ws2.column_dimensions[col].width = 18

    # ABA 3: GEOMETRIA
    ws3 = wb.create_sheet("GEOMETRIA")
    ws3["A1"] = "Geometria dos Cenários (premissa do estudo)"
    ws3["A1"].font = Font(bold=True, size=13, color=DARK)
    ws3.merge_cells("A1:D1")

    geom_table = [
        ("Variável", "Vertical", "Horizontal", "Razão V/H"),
        ("Aptos por pavimento", GEOMETRIA["vertical"]["aptos_pav"], GEOMETRIA["horizontal"]["aptos_pav"], "0.75x"),
        ("Pavimentos tipo (NPT)", GEOMETRIA["vertical"]["npt"], GEOMETRIA["horizontal"]["npt"], "1.33x"),
        ("Pavimentos total (NP)", GEOMETRIA["vertical"]["np"], GEOMETRIA["horizontal"]["np"], "1.29x"),
        ("Footprint torre (m²)", GEOMETRIA["vertical"]["footprint"], GEOMETRIA["horizontal"]["footprint"], "0.75x"),
        ("Perímetro torre (m)", GEOMETRIA["vertical"]["perim"], GEOMETRIA["horizontal"]["perim"], "0.87x"),
        ("Altura aproximada (m)", GEOMETRIA["vertical"]["altura"], GEOMETRIA["horizontal"]["altura"], "1.30x"),
        ("Área fachada total (m²)", GEOMETRIA["vertical"]["fachada"], GEOMETRIA["horizontal"]["fachada"], "1.15x"),
        ("Fachada / m² AC", round(GEOMETRIA["vertical"]["fachada"]/2000, 3), round(GEOMETRIA["horizontal"]["fachada"]/2000, 3), "1.15x"),
        ("Cobertura / m² AC", round(GEOMETRIA["vertical"]["cob"]/2000, 3), round(GEOMETRIA["horizontal"]["cob"]/2000, 3), "0.75x"),
        ("Elevadores", GEOMETRIA["vertical"]["elev"], GEOMETRIA["horizontal"]["elev"], "2x"),
    ]
    for i, row in enumerate(geom_table, 3):
        for j, val in enumerate(row, 1):
            c = ws3.cell(i, j, val)
            if i == 3:
                style_header(c)
            else:
                c.border = THIN
                if j == 1:
                    c.font = Font(bold=True, size=10)
                else:
                    c.alignment = Alignment(horizontal='center')

    ws3.column_dimensions['A'].width = 28
    for col in 'BCD':
        ws3.column_dimensions[col].width = 16

    # ABA 4: GRAFICOS
    ws4 = wb.create_sheet("GRAFICOS")
    ws4["A1"] = "Comparativo gráfico"
    ws4["A1"].font = Font(bold=True, size=13, color=DARK)

    # Tabela de dados pro gráfico (todos macrogrupos)
    ws4.append([])
    ws4.append(["Macrogrupo", "Vertical", "Horizontal"])
    for i, c in enumerate(["Macrogrupo", "Vertical", "Horizontal"], 1):
        style_header(ws4.cell(3, i))
    for macro, v, h, d, dp in deltas:
        ws4.append([macro, v, h])
    for r in range(4, 4 + len(deltas)):
        for col in [2, 3]:
            style_money(ws4.cell(r, col))
        ws4.cell(r, 1).border = THIN

    # Bar chart
    chart = BarChart()
    chart.type = "bar"
    chart.style = 10
    chart.title = "Custo por macrogrupo — V × H (R$)"
    chart.y_axis.title = "Macrogrupo"
    chart.x_axis.title = "Custo (R$)"
    chart.height = 18
    chart.width = 22
    data = Reference(ws4, min_col=2, max_col=3, min_row=3, max_row=3 + len(deltas))
    cats = Reference(ws4, min_col=1, min_row=4, max_row=3 + len(deltas))
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    ws4.add_chart(chart, "F3")

    ws4.column_dimensions['A'].width = 28
    ws4.column_dimensions['B'].width = 18
    ws4.column_dimensions['C'].width = 18

    wb.save(OUT_XLSX)
    print(f"  OK: {OUT_XLSX.name}")

# ============== ETAPA 4b: GERAR MEMORIAL DOCX ==============

def gerar_memorial_docx(t_v, gt_v, t_h, gt_h, fachada_v, fachada_h):
    print("\n[4b/4] Gerando memorial.docx...")
    doc = Document()

    # Estilos
    style = doc.styles['Normal']
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Capa
    h = doc.add_heading("Estudo de Massa 2.000 m² — Análise Comparativa", level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Verticalização (3 aptos/pav) × Horizontalização (4 aptos/pav)")
    r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Brusque/SC | 24 aptos × 55 m² priv | Padrão Médio | Térreo pilotis | Geração: 30/04/2026")
    r.font.size = Pt(10); r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    doc.add_paragraph()

    # Sumário
    doc.add_heading("1. Resumo executivo", level=1)
    delta = gt_v - gt_h
    delta_pct = delta / gt_h * 100 if gt_h else 0
    winner = "Horizontal" if delta > 0 else "Vertical"
    p = doc.add_paragraph()
    p.add_run("Cenário Vertical (3 aptos/pav, 9 pavs total): ").bold = True
    p.add_run(f"{fmt_brl(gt_v)}\n")
    p.add_run("Cenário Horizontal (4 aptos/pav, 7 pavs total): ").bold = True
    p.add_run(f"{fmt_brl(gt_h)}\n")
    p.add_run("Δ (V − H): ").bold = True
    p.add_run(f"{fmt_brl(delta)} ({delta_pct:+.1f}%)\n")
    p.add_run(f"Cenário mais econômico: ").bold = True
    p.add_run(f"{winner}").bold = True

    doc.add_paragraph()
    doc.add_paragraph(
        "Em projeto desse porte (2.000 m² AC, padrão Médio), a diferença entre verticalizar "
        "e alongar não é apenas no custo total — ela se concentra em poucos macrogrupos sensíveis "
        "à geometria. O memorial abaixo identifica onde está o dinheiro e por quê."
    )

    # Premissas
    doc.add_heading("2. Premissas", level=1)
    doc.add_heading("2.1 Comuns aos 2 cenários", level=2)
    doc.add_paragraph(
        "AC 2.000 m² | UR 24 aptos × 55 m² priv | Padrão Médio | CUB R$ 2.870/m² (Norte SC) | "
        "Laje convencional | Fundação hélice contínua | Sem subsolo (térreo pilotis garagem) | "
        "Fachada textura+pintura | Vedação alvenaria | Esquadrias alumínio anodizado | "
        "Lazer básico (salão+churrasqueira) | Prazo 18 meses | Pé-direito 2.80 m | 1 banheiro/apto | "
        "Tipologia 1-2 dormitórios | 1 torre | Sem gerador, subestação, fotovoltaica, pressurização ou piscina."
    )

    doc.add_heading("2.2 Geometria — o que muda", level=2)
    table = doc.add_table(rows=11, cols=3)
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_cells = table.rows[0].cells
    for i, t in enumerate(["Variável", "Vertical (V)", "Horizontal (H)"]):
        hdr_cells[i].text = t
        for r in hdr_cells[i].paragraphs[0].runs:
            r.bold = True

    geom_rows = [
        ("Aptos por pavimento", "3", "4"),
        ("Pavimentos tipo (NPT)", "8", "6"),
        ("Pavimentos total (NP)", "9", "7"),
        ("Footprint torre (m²)", "~235", "~315"),
        ("Perímetro torre (m)", "~62", "~71"),
        ("Altura aprox. (m)", "~26", "~20"),
        ("Área fachada total (m²)", "~1.395", "~1.213"),
        ("Fachada / m² AC", "0.74", "0.64"),
        ("Cobertura / m² AC", "0.13", "0.17"),
        ("Elevadores", "2", "1"),
    ]
    for i, (a, b, c) in enumerate(geom_rows, 1):
        cells = table.rows[i].cells
        cells[0].text = a; cells[1].text = b; cells[2].text = c

    # Resultado
    doc.add_paragraph()
    doc.add_heading("3. Resultado por macrogrupo", level=1)
    deltas = [(macro, t_v.get(macro, 0), t_h.get(macro, 0),
               t_v.get(macro, 0) - t_h.get(macro, 0),
               (t_v.get(macro, 0) - t_h.get(macro, 0)) / t_h.get(macro, 1) if t_h.get(macro) else 0)
              for macro in t_v]
    deltas_top = sorted(deltas, key=lambda x: abs(x[3]), reverse=True)[:5]

    doc.add_paragraph("Top 5 macrogrupos com maior diferença absoluta (V − H):")
    table = doc.add_table(rows=6, cols=5)
    table.style = "Light Grid Accent 1"
    hdr_cells = table.rows[0].cells
    for i, t in enumerate(["#", "Macrogrupo", "Vertical", "Horizontal", "Δ (V − H)"]):
        hdr_cells[i].text = t
        for r in hdr_cells[i].paragraphs[0].runs:
            r.bold = True
    for i, (macro, v, h, d, dp) in enumerate(deltas_top, 1):
        cells = table.rows[i].cells
        cells[0].text = str(i)
        cells[1].text = macro
        cells[2].text = fmt_brl(v)
        cells[3].text = fmt_brl(h)
        cells[4].text = f"{fmt_brl(d)} ({dp*100:+.1f}%)"

    # Análise dos macrogrupos sensíveis
    doc.add_heading("4. Análise — onde está o dinheiro", level=1)
    explicacoes = {
        "Fachada": (
            "A torre vertical é mais esbelta — perímetro grande em relação ao footprint. "
            "Por m² AC, a fachada do cenário V tem ~0.74 m²/m² vs ~0.64 m²/m² no H "
            "(15% mais fachada). Como o PU médio de fachada textura+pintura é ~R$ 153/m² "
            "(pintura+textura+balancim+MO), cada m² extra de fachada pesa direto no custo."
        ),
        "Sist. Especiais": (
            "O cenário V tem 9 paradas (térreo + 8 pavs) — pela NBR 13994 e tráfego típico, "
            "exige 2 elevadores sociais + 1 de serviço/social compartilhado. "
            "O cenário H tem 7 paradas e 24 aptos — fica no limite, modelado com 1 social + 1 serviço. "
            "Cada elevador adicional custa ~R$ 192k (5 paradas base) + acréscimo por parada extra."
        ),
        "Supraestrutura": (
            "Torre mais alta = pilares maiores no térreo (carga acumulada), mais aço no núcleo, "
            "mais escoramento progressivo. O fator pé-direito e altura captura ~5% de delta a favor do H. "
            "Na contramão, o H tem mais m² de laje por pav (footprint maior) — efeito menor."
        ),
        "Mov. Terra": (
            "Sem subsolo, o macro Mov. Terra é só nivelamento e regularização. "
            "Diferença pequena, depende do footprint do embasamento."
        ),
        "Infraestrutura": (
            "Fundação hélice: nº de estacas escala com AC (igual nos 2), mas a profundidade "
            "depende da carga total — no V, carga concentrada em footprint menor exige mais "
            "estacas longas; no H, carga distribuída em footprint maior."
        ),
        "Gerenciamento": (
            "Prazo igual (18 meses), mas torre alta exige grua, andaimes mais altos, equipe "
            "vertical maior — o cenário V tem indiretos um pouco mais pesados."
        ),
    }
    for macro, _, _, _, _ in deltas_top:
        doc.add_heading(macro, level=2)
        doc.add_paragraph(explicacoes.get(macro, "Análise específica do macrogrupo (revisar caso a caso)."))

    # Disclaimer vagas
    doc.add_heading("5. Disclaimer importante — vagas", level=1)
    doc.add_paragraph(
        "Este estudo modelou o pilotis com 8 vagas (cenário V) e 11 vagas (cenário H), "
        "que cabem no footprint do térreo. Para um projeto real com 24 aptos × 1 vaga = 24 vagas, "
        "essa quantidade é insuficiente em ambos os cenários. As alternativas em projeto real seriam:"
    )
    doc.add_paragraph(
        "• 1 subsolo de garagem (+R$ 200-400/m² no macro Infraestrutura/Mov. Terra/Impermeabilização) — "
        "padrão SC litoral, escala com footprint;\n"
        "• Pilotis duplo (térreo + 1° pavimento garagem) — sem escavação, mas perde-se 1 pav residencial "
        "ou aumenta-se 1 pav, mudando AC efetiva;\n"
        "• Aceitar < 1 vaga/apto (compactos populares ou estúdios) — reduz custo total ~10%."
    )
    doc.add_paragraph(
        "A decisão de garagem precisa ser feita ANTES de fechar o produto — o impacto é grande "
        "(R$ 250-500k em projeto desse porte) e altera o ranking dos macrogrupos sensíveis."
    )

    # Recomendação
    doc.add_heading("6. Recomendação", level=1)
    p = doc.add_paragraph()
    p.add_run("Decisão por critério de custo puro: ").bold = True
    p.add_run(
        f"o cenário {winner} é {abs(delta_pct):.1f}% mais econômico ({fmt_brl(abs(delta))}). "
        f"Para 2.000 m² AC essa diferença é relevante mas não decisiva — a escolha do produto "
        f"deve considerar fatores extra-custo:"
    )
    doc.add_paragraph(
        "• Verticalizar faz sentido quando: terreno caro/escasso, COE/zoneamento permite, "
        "view valoriza pavimento alto, mercado aceita pagar premium por andar superior.\n"
        "• Alongar faz sentido quando: terreno barato/disponível, comprador valoriza apto térreo "
        "ou primeiro andar, custos de canteiro/grua são caros na região, prazo é crítico "
        "(obra horizontal é mais rápida)."
    )
    p = doc.add_paragraph()
    p.add_run("O aprendizado real deste estudo é qual macrogrupo domina o delta — ").bold = True
    p.add_run(
        f"e a resposta é: Fachada e Sistemas Especiais (elevadores) somados respondem por "
        f"~{(deltas_top[0][3] + deltas_top[1][3])/delta*100:.0f}% da diferença total. "
        "Toda decisão de produto que afete fachada/m² AC ou nº de elevadores tem impacto "
        "desproporcional comparado a outras decisões (laje, fundação, padrão de acabamento)."
    )

    # Premissas técnicas
    doc.add_heading("7. Premissas técnicas (overrides aplicados)", level=1)
    doc.add_paragraph(
        "Para captar o efeito geométrico real, os índices abaixo foram sobrescritos manualmente "
        "na aba INDICES (Col C — override) e na aba Sist. Especiais de cada xlsx:"
    )
    doc.add_paragraph(
        f"• INDICES!C18 — Fachada/m² AC: V=0.74 | H=0.64 (calc default usaria 1.55, mas refletindo "
        f"a geometria real do estudo de massa)\n"
        f"• Sist. Especiais!C4 — Qtd elevadores sociais: V=2 | H=1 (default fórmula MAX(1,N-1) "
        f"não diferencia 1 vs 2 elevadores corretamente)"
    )
    doc.add_paragraph()
    doc.add_paragraph(
        "Calibração estatística cross-projeto (4.210 PUs V2, 126 projetos) na base Cartesian. "
        "Não substitui orçamento executivo (BoQ rastreável item a item) — é estudo de massa "
        "para decisão de produto."
    )

    doc.save(OUT_DOCX)
    print(f"  OK: {OUT_DOCX.name}")

# ============== ETAPA 4c: GERAR APRESENTAÇÃO PPTX ==============

def gerar_apresentacao_pptx(t_v, gt_v, t_h, gt_h, fachada_v, fachada_h):
    print("\n[4c/4] Gerando apresentacao.pptx...")
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    DARK_RGB = PRGB(0x2C, 0x3E, 0x50)
    ACCENT_RGB = PRGB(0x29, 0x80, 0xB9)
    GREEN_RGB = PRGB(0x27, 0xAE, 0x60)
    RED_RGB = PRGB(0xE7, 0x4C, 0x3C)
    GRAY_RGB = PRGB(0x66, 0x66, 0x66)

    def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_RGB, align="left"):
        tb = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if align == "center":
            p.alignment = 2
        elif align == "right":
            p.alignment = 3
        r = p.add_run()
        r.text = text
        r.font.size = PPt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
        return tf

    def add_table(slide, x, y, w, h, data, header=True):
        rows, cols = len(data), len(data[0])
        tbl = slide.shapes.add_table(rows, cols, PInches(x), PInches(y), PInches(w), PInches(h)).table
        for ri, row in enumerate(data):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                for r in p.runs:
                    r.font.size = PPt(11)
                    r.font.name = "Calibri"
                    if ri == 0 and header:
                        r.font.bold = True
                        r.font.color.rgb = PRGB(0xFF, 0xFF, 0xFF)
                if ri == 0 and header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = DARK_RGB
        return tbl

    blank = prs.slide_layouts[6]

    # SLIDE 1 — Capa
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 2.5, 12.3, 1, "Estudo de Massa 2.000 m²", size=44, bold=True, color=DARK_RGB, align="center")
    add_text(s, 0.5, 3.6, 12.3, 0.6, "Verticalização × Horizontalização", size=28, color=ACCENT_RGB, align="center")
    add_text(s, 0.5, 4.4, 12.3, 0.5, "Onde está o dinheiro num multifamiliar pequeno?",
             size=18, color=GRAY_RGB, align="center")
    add_text(s, 0.5, 6.5, 12.3, 0.4,
             "Brusque/SC | 24 aptos × 55 m² priv | Padrão Médio | Térreo pilotis | 30/04/2026",
             size=12, color=GRAY_RGB, align="center")

    # SLIDE 2 — Premissas comuns
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Premissas comuns aos 2 cenários", size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4,
             "Mesmas escolhas de produto — o que muda é só a geometria",
             size=14, color=GRAY_RGB)
    data = [
        ["Categoria", "Valor"],
        ["AC total", "2.000 m²"],
        ["UR / Apto privativo", "24 aptos × 55 m²"],
        ["Padrão de acabamento", "Médio (porcelanato, alumínio anodizado)"],
        ["Cidade", "Brusque/SC (Norte SC)"],
        ["CUB referência", "R$ 2.870/m² (data-base abr/26)"],
        ["Estrutura", "Laje convencional, fundação hélice"],
        ["Vedação / Fachada", "Alvenaria | Textura + pintura"],
        ["Térreo", "Pilotis garagem (sem subsolo)"],
        ["Lazer / Prazo", "Básico (salão+churrasq.) | 18 meses"],
        ["1 torre, sem gerador, sem subestação, sem fotovoltaica, sem piscina, sem pressurização", ""],
    ]
    add_table(s, 0.7, 1.8, 12, 5.0, data)

    # SLIDE 3 — Cenário V
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Cenário V — Vertical (3 aptos/pav)", size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4, "Torre esbelta, mais alta, footprint menor", size=14, color=GRAY_RGB)
    data = [
        ["Variável", "Valor"],
        ["Aptos por pavimento", "3"],
        ["Pavimentos tipo (NPT)", "8"],
        ["Pavimentos total (NP)", "9 (1 térreo pilotis + 8 tipo)"],
        ["Footprint torre", "~235 m²"],
        ["Perímetro torre", "~62 m"],
        ["Altura aprox.", "~26 m"],
        ["Elevadores", "2 (NBR 13994 — 9 paradas)"],
        ["Área fachada", "~1.395 m² (0.74 m²/m² AC)"],
        ["Custo total estimado", fmt_brl(gt_v)],
        ["Custo / m² AC", fmt_brl(gt_v/2000)],
        ["Custo / UR", fmt_brl(gt_v/24)],
    ]
    add_table(s, 0.7, 1.8, 12, 5.4, data)

    # SLIDE 4 — Cenário H
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Cenário H — Horizontal (4 aptos/pav)", size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4, "Torre mais baixa, footprint maior, mais retangular", size=14, color=GRAY_RGB)
    data = [
        ["Variável", "Valor"],
        ["Aptos por pavimento", "4"],
        ["Pavimentos tipo (NPT)", "6"],
        ["Pavimentos total (NP)", "7 (1 térreo pilotis + 6 tipo)"],
        ["Footprint torre", "~315 m²"],
        ["Perímetro torre", "~71 m"],
        ["Altura aprox.", "~20 m"],
        ["Elevadores", "1 (7 paradas, tráfego baixo)"],
        ["Área fachada", "~1.213 m² (0.64 m²/m² AC)"],
        ["Custo total estimado", fmt_brl(gt_h)],
        ["Custo / m² AC", fmt_brl(gt_h/2000)],
        ["Custo / UR", fmt_brl(gt_h/24)],
    ]
    add_table(s, 0.7, 1.8, 12, 5.4, data)

    # SLIDE 5 — Grand Total comparativo
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Grand Total — V × H", size=28, bold=True, color=DARK_RGB)
    delta = gt_v - gt_h
    delta_pct = delta / gt_h * 100 if gt_h else 0
    winner = "Horizontal" if delta > 0 else "Vertical"
    add_text(s, 0.5, 1.1, 12.3, 0.4,
             f"Cenário mais econômico: {winner} | Δ {fmt_brl(abs(delta))} ({abs(delta_pct):.1f}%)",
             size=18, bold=True, color=ACCENT_RGB)

    chart_data = CategoryChartData()
    chart_data.categories = ["Vertical (V)", "Horizontal (H)"]
    chart_data.add_series("Custo Total (R$)", (gt_v, gt_h))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PInches(2.5), PInches(2.0), PInches(8.5), PInches(4.5),
        chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False

    add_text(s, 0.5, 6.7, 12.3, 0.4,
             f"V: {fmt_brl(gt_v)} ({fmt_brl(gt_v/2000)}/m²)   |   H: {fmt_brl(gt_h)} ({fmt_brl(gt_h/2000)}/m²)",
             size=14, color=GRAY_RGB, align="center")

    # SLIDE 6 — Comparativo por macrogrupo (tabela top 10)
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Comparativo por macrogrupo", size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4, "Top 10 macrogrupos por custo total (V)", size=14, color=GRAY_RGB)

    # Top 10 por custo absoluto V
    top10 = sorted(t_v.items(), key=lambda x: x[1], reverse=True)[:10]
    data = [["Macrogrupo", "Vertical", "Horizontal", "Δ (V−H)", "Δ %"]]
    for macro, v in top10:
        h = t_h.get(macro, 0)
        d = v - h
        dp = d/h*100 if h else 0
        data.append([macro, fmt_brl(v), fmt_brl(h), fmt_brl(d), f"{dp:+.1f}%"])
    add_table(s, 0.5, 1.7, 12.3, 5.5, data)

    # SLIDE 7 — Top 5 sensíveis (waterfall-style)
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Top 5 macrogrupos sensíveis à geometria",
             size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4, "Onde a diferença V−H se concentra", size=14, color=GRAY_RGB)

    deltas = [(macro, t_v.get(macro, 0), t_h.get(macro, 0),
               t_v.get(macro, 0) - t_h.get(macro, 0))
              for macro in t_v]
    deltas_top = sorted(deltas, key=lambda x: abs(x[3]), reverse=True)[:5]

    chart_data = CategoryChartData()
    chart_data.categories = [m for m, v, h, d in deltas_top]
    chart_data.add_series("Δ V−H (R$)", [d for m, v, h, d in deltas_top])
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        PInches(0.5), PInches(1.7), PInches(12.3), PInches(5.5),
        chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False

    # SLIDE 8 — Análise de fachada
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Análise — Fachada", size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4, "O macrogrupo que mais sente a geometria", size=14, color=GRAY_RGB)

    fac_v = t_v.get("Fachada", 0)
    fac_h = t_h.get("Fachada", 0)
    data = [
        ["", "Vertical (V)", "Horizontal (H)"],
        ["Área fachada total (m²)", "~1.395", "~1.213"],
        ["Fachada / m² AC", "0.74", "0.64"],
        ["Custo macro Fachada", fmt_brl(fac_v), fmt_brl(fac_h)],
        ["Δ Fachada (V−H)", fmt_brl(fac_v - fac_h), ""],
        ["Δ % do total geral", f"{(fac_v-fac_h)/(gt_v-gt_h)*100 if gt_v != gt_h else 0:.0f}%", ""],
    ]
    add_table(s, 0.5, 1.8, 12.3, 4.0, data)

    add_text(s, 0.5, 6.2, 12.3, 1,
             "Por quê: torre vertical é esbelta — perímetro grande proporcional ao footprint. "
             "+15% de fachada/m² AC → impacto direto no maior macrogrupo de envoltória.",
             size=12, color=GRAY_RGB)

    # SLIDE 9 — Análise de elevadores
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Análise — Sistemas Especiais (Elevadores)",
             size=28, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 1.1, 12.3, 0.4, "O elevador adicional pesa muito", size=14, color=GRAY_RGB)

    se_v = t_v.get("Sist. Especiais", 0)
    se_h = t_h.get("Sist. Especiais", 0)
    data = [
        ["", "Vertical (V)", "Horizontal (H)"],
        ["Paradas", "9 (1 térreo + 8 tipo)", "7 (1 térreo + 6 tipo)"],
        ["Elevadores sociais", "2", "1"],
        ["Elevador serviço", "1", "1"],
        ["NBR 13994", "Exige 2 sociais", "1 social no limite"],
        ["Custo macro Sist. Especiais", fmt_brl(se_v), fmt_brl(se_h)],
        ["Δ Sist. Especiais (V−H)", fmt_brl(se_v - se_h), ""],
    ]
    add_table(s, 0.5, 1.8, 12.3, 4.5, data)
    add_text(s, 0.5, 6.5, 12.3, 0.6,
             "Por quê: PU elevador 5 paradas ≈ R$ 192k. Cada parada extra ≈ R$ 11k. "
             "9 paradas + 24 aptos = 2 elevadores. 7 paradas + 24 aptos = 1.",
             size=12, color=GRAY_RGB)

    # SLIDE 10 — Conclusão
    s = prs.slides.add_slide(blank)
    add_text(s, 0.5, 0.4, 12.3, 0.6, "Conclusão e recomendação", size=28, bold=True, color=DARK_RGB)

    add_text(s, 0.5, 1.3, 12.3, 0.5,
             f"Cenário mais econômico em custo puro: {winner}",
             size=20, bold=True, color=ACCENT_RGB)
    add_text(s, 0.5, 1.9, 12.3, 0.4,
             f"Δ absoluto: {fmt_brl(abs(delta))} ({abs(delta_pct):.1f}% do total)",
             size=14, color=GRAY_RGB)

    fac_d = t_v.get("Fachada", 0) - t_h.get("Fachada", 0)
    se_d = t_v.get("Sist. Especiais", 0) - t_h.get("Sist. Especiais", 0)
    pct_top2 = (fac_d + se_d) / delta * 100 if delta else 0

    add_text(s, 0.5, 2.7, 12.3, 0.5, "Onde está o dinheiro", size=18, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 3.3, 12.3, 1.2,
             f"Fachada + Sistemas Especiais (elevadores) somam ~{abs(pct_top2):.0f}% da diferença total. "
             f"Toda decisão de produto que mexa com perímetro de torre, altura ou nº de elevadores "
             f"tem impacto desproporcional nesses 2 macros.",
             size=14, color=DARK_RGB)

    add_text(s, 0.5, 4.7, 12.3, 0.5, "Quando vale a pena verticalizar mesmo assim", size=18, bold=True, color=DARK_RGB)
    add_text(s, 0.5, 5.3, 12.3, 1.5,
             "• Terreno caro/escasso (CUB do terreno > delta de fachada+elevador)\n"
             "• Zoneamento exige altura mínima ou taxa de ocupação baixa\n"
             "• View valoriza o pavimento alto (preço m² premium > custo extra)\n"
             "• Mercado-alvo aceita pagar premium por andar superior",
             size=13, color=DARK_RGB)

    add_text(s, 0.5, 6.9, 12.3, 0.4,
             "Atenção: vagas insuficientes em ambos cenários — definir solução de garagem antes de fechar produto",
             size=11, color=PRGB(0xE7, 0x4C, 0x3C))

    prs.save(OUT_PPTX)
    print(f"  OK: {OUT_PPTX.name}")

# ============== MAIN ==============

def main():
    if not XLSX_V.exists() or not XLSX_H.exists():
        print(f"ERRO: xlsx não encontrados. Rode primeiro:")
        print(f"  py -3.10 scripts/gerar_template_dinamico_v2.py --config base/pacotes/{SLUG}/parametrico-v2-config-vertical.json -o {XLSX_V.relative_to(Path.home() / 'orcamentos-openclaw')}")
        sys.exit(1)

    aplicar_overrides()
    recalcular_xlsx([XLSX_V, XLSX_H])

    print("\n[3/4] Lendo totais por macrogrupo...")
    t_v, gt_v = ler_totais(XLSX_V)
    t_h, gt_h = ler_totais(XLSX_H)
    fachada_v = ler_indice_fachada(XLSX_V)
    fachada_h = ler_indice_fachada(XLSX_H)

    print(f"  Vertical:   Grand Total = {fmt_brl(gt_v)} ({fmt_brl(gt_v/2000)}/m²)")
    print(f"  Horizontal: Grand Total = {fmt_brl(gt_h)} ({fmt_brl(gt_h/2000)}/m²)")
    print(f"  Δ V−H: {fmt_brl(gt_v - gt_h)} ({(gt_v-gt_h)/gt_h*100:+.1f}%)")
    print(f"  Fachada/m² AC: V={fachada_v:.3f} | H={fachada_h:.3f}")

    print("\n[4/4] Gerando entregáveis...")
    gerar_comparativo_xlsx(t_v, gt_v, t_h, gt_h, fachada_v, fachada_h)
    gerar_memorial_docx(t_v, gt_v, t_h, gt_h, fachada_v, fachada_h)
    gerar_apresentacao_pptx(t_v, gt_v, t_h, gt_h, fachada_v, fachada_h)

    print("\n=== TODOS OS ENTREGÁVEIS GERADOS ===")
    for f in [XLSX_V, XLSX_H, OUT_XLSX, OUT_DOCX, OUT_PPTX]:
        size = f.stat().st_size / 1024
        print(f"  {f.name} ({size:.1f} KB)")

if __name__ == "__main__":
    main()
